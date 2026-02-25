import os
# Set tokenizers parallelism to avoid warnings when forking
os.environ["TOKENIZERS_PARALLELISM"] = "false"

# Load environment variables from .env file (OS env vars take priority)
# override=False ensures OS environment variables are not overwritten by .env file
try:
    from dotenv import load_dotenv
    # Check if .env file exists before loading
    env_file_path = os.path.join(os.path.dirname(__file__), ".env")
    if os.path.exists(env_file_path):
        load_dotenv(override=False)
        # Only print if DEBUG mode is enabled (check after loading .env)
        # Note: We need to check again after load_dotenv since DEBUG might be in .env
        if os.getenv("DEBUG", "false").lower() == "true":
            print("Environment variables loaded (OS env vars take priority, .env file as fallback)", flush=True)
except ImportError:
    print("Warning: python-dotenv not installed. Install with: pip install python-dotenv", flush=True)


# Sanitize OMP_NUM_THREADS before importing torch to prevent libgomp errors
if "OMP_NUM_THREADS" in os.environ:
    omp_val = os.environ["OMP_NUM_THREADS"]
    # If empty or not a valid positive integer, unset it to let torch/OpenMP use defaults
    if not omp_val or not omp_val.isdigit() or int(omp_val) <= 0:
        del os.environ["OMP_NUM_THREADS"]

import gradio as gr

import gc

from transformers import AutoTokenizer, AutoModelForSequenceClassification
import torch
import re
from tokenizers import normalizers
from tokenizers.normalizers import Sequence, Replace, Strip, NFKC
from tokenizers import Regex
import matplotlib.pyplot as plt
import html
import pymupdf
import warnings
from docx import Document
from pptx import Presentation
import pandas as pd
import shutil
from pathlib import Path
import time
import language_tool_python
from langdetect import detect, DetectorFactory
import nltk
from nltk.tokenize import sent_tokenize, word_tokenize

from nltk.corpus import stopwords
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import difflib
from collections import Counter
import requests
import json
import math

# Try to import and load spaCy for sentence splitting
try:
    import spacy
    # Load multilingual UD-based sentencizer model
    try:
        nlp = spacy.load("xx_sent_ud_sm")
        SPACY_AVAILABLE = True
        print("spaCy multilingual sentence tokenizer loaded successfully.", flush=True)
    except OSError:
        # Model not found, try to download it or use fallback
        print("Warning: spaCy model 'xx_sent_ud_sm' not found. Install with: python -m spacy download xx_sent_ud_sm", flush=True)
        try:
            # Try loading English model as fallback
            nlp = spacy.load("en_core_web_sm")
            SPACY_AVAILABLE = True
            print("Using spaCy English model as fallback.", flush=True)
        except OSError:
            nlp = None
            SPACY_AVAILABLE = False
            print("Warning: spaCy models not available. Using NLTK/Regex fallback.", flush=True)
except ImportError:
    spacy = None
    nlp = None
    SPACY_AVAILABLE = False
    print("Warning: spaCy not installed. Using NLTK/Regex fallback for sentence splitting.", flush=True)

# Suppress warnings
warnings.filterwarnings('ignore', category=UserWarning)
warnings.filterwarnings('ignore', category=FutureWarning)
import logging
logging.getLogger("pymupdf").setLevel(logging.ERROR)

# Suppress transformers warnings about uninitialized weights
os.environ["TRANSFORMERS_VERBOSITY"] = "error"
logging.getLogger("transformers").setLevel(logging.ERROR)
logging.getLogger("transformers.modeling_utils").setLevel(logging.ERROR)
warnings.filterwarnings('ignore', message='.*Some weights of.*were not initialized.*')
warnings.filterwarnings('ignore', message='.*You should probably TRAIN this model.*')
warnings.filterwarnings('ignore', message='.*were not initialized from the model checkpoint.*')

# Set device - prioritize CUDA, then MPS (Apple Silicon), then CPU
if torch.cuda.is_available():
    device = torch.device('cuda')
    print(f"CUDA available! Using GPU: {torch.cuda.get_device_name(0)}", flush=True)
elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
    device = torch.device('mps')
    print(f"MPS (Apple Silicon GPU) available! Using Metal Performance Shaders.", flush=True)
else:
    device = torch.device('cpu')
    print(f"Using CPU for inference.", flush=True)

# Performance optimizations
if torch.cuda.is_available():
    # Enable cuDNN benchmarking for faster inference (finds optimal algorithms)
    torch.backends.cudnn.benchmark = True
    torch.backends.cudnn.deterministic = False
    print(f"CUDA version: {torch.version.cuda}", flush=True)
elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
    # MPS optimizations for Apple Silicon
    print(f"MPS optimizations enabled for Apple Silicon.", flush=True)
else:
    # Optimize CPU performance
    # Use all available CPU cores if OMP_NUM_THREADS is not set
    num_threads_str = os.environ.get('OMP_NUM_THREADS')
    
    # Try to parse valid thread count
    valid_threads = False
    if num_threads_str and num_threads_str.isdigit():
        num_threads = int(num_threads_str)
        if num_threads > 0:
            torch.set_num_threads(num_threads)
            torch.set_num_interop_threads(num_threads)
            print(f"Using CPU with {num_threads} threads", flush=True)
            valid_threads = True
            
    if not valid_threads:
        # Use all available cores
        import multiprocessing
        num_cores = multiprocessing.cpu_count()
        torch.set_num_threads(num_cores)
        torch.set_num_interop_threads(num_cores)
        print(f"Using CPU with all available cores ({num_cores} threads)", flush=True)
    
    # Enable optimizations for CPU inference
    torch.set_float32_matmul_precision('medium')  # Use faster but less precise matmul

# Ensure uploads directory exists
UPLOADS_DIR = Path("uploads")
UPLOADS_DIR.mkdir(exist_ok=True)

def cleanup_memory():
    """Force garbage collection and empty GPU cache"""
    try:
        if torch.cuda.is_available():
            torch.cuda.empty_cache()
        elif hasattr(torch.backends, 'mps') and torch.backends.mps.is_available():
            try:
                torch.mps.empty_cache()
            except:
                pass
    except:
        pass
    gc.collect()


# --- Language Detection ---
def detect_language(text):
    """
    Detect if the text is in English.
    Returns True if English (any variant), False otherwise.
    
    Supports all English variants:
    - en-US (US English)
    - en-GB (British English)
    - en-CA (Canadian English)
    - en-AU (Australian English)
    - en-NZ (New Zealand English)
    - en-ZA (South African English)
    
    Note: langdetect returns 'en' for all English variants (ISO 639-1 code).
    For non-English languages, it returns their ISO 639-1 codes (e.g., 'ur' for Urdu, 'ar' for Arabic, etc.)
    """
    if not text or not text.strip():
        return True  # Empty text is considered valid (will be caught by other validations)
    
    try:
        # Set seed for consistent results
        DetectorFactory.seed = 0
        
        # Use a sample of text for detection (first 2000 chars should be enough)
        # This helps with performance and langdetect works better with reasonable length text
        text_sample = text.strip()[:2000] if len(text.strip()) > 2000 else text.strip()
        
        detected_lang = detect(text_sample)
        
        # langdetect returns ISO 639-1 codes (2 letters)
        # For all English variants, it returns 'en'
        # For non-English: 'ur' (Urdu), 'ar' (Arabic), 'hi' (Hindi), 'zh' (Chinese), etc.
        is_english = detected_lang == 'en' or detected_lang.startswith('en')
        
        return is_english
    except Exception as e:
        # If detection fails, assume non-English to be safe
        return False

def check_language_and_return_error(text):
    """
    Check if text is English. If not, return error message.
    Returns (is_valid, error_message)
    """
    if not text or not text.strip():
        return True, None  # Empty text will be caught by other validations
    
    if not detect_language(text):
        error_msg = '<div style="color: red; font-weight: bold;">Error: Only English language is supported. Please provide text in English.</div>'
        return False, error_msg
    
    return True, None

# Cache to track files that have already been saved (to avoid duplicates)
_saved_files_cache = {}

def save_uploaded_file(file_path: str) -> str:
    """Copy uploaded file to uploads directory and return the new path"""
    if not file_path or not os.path.exists(file_path):
        return file_path
    
    # Check if file is already in uploads directory
    try:
        file_path_abs = Path(file_path).resolve()
        uploads_abs = UPLOADS_DIR.resolve()
        # If file is already in uploads directory, return it as-is
        if uploads_abs in file_path_abs.parents or file_path_abs.parent == uploads_abs:
            return str(file_path_abs)
    except:
        pass
    
    # Get the filename from the original path
    filename = os.path.basename(file_path)
    
    # Check cache to see if we've already saved this file (by filename)
    # This handles cases where Gradio provides different temp paths for the same file
    if filename in _saved_files_cache:
        cached_path = _saved_files_cache[filename]
        # Verify the cached file still exists
        if os.path.exists(cached_path):
            return cached_path
    
    # Create destination path in uploads directory
    dest_path = UPLOADS_DIR / filename
    
    # Check if file with same name exists - replace it if it does
    if dest_path.exists():
        try:
            source_size = os.path.getsize(file_path)
            dest_size = os.path.getsize(dest_path)
            # If files have the same size, assume it's the same file and return existing path
            if source_size == dest_size:
                _saved_files_cache[filename] = str(dest_path)
                return str(dest_path)
        except:
            pass
        
        # File exists but sizes differ (or comparison failed), replace the existing file
        # Remove the old file before copying the new one
        try:
            os.remove(dest_path)
        except:
            pass
    
    # Copy file to uploads directory (will overwrite if file exists)
    shutil.copy2(file_path, dest_path)
    
    saved_path = str(dest_path)
    # Cache the mapping from filename to saved path
    _saved_files_cache[filename] = saved_path
    
    return saved_path

# --- Model and Tokenizer Setup ---
# Local model directory
LOCAL_MODEL_DIR = "model/ModernBERT-base"
model1_path = "model/modernbert.bin"
model2_path = "model/Model_groups_3class_seed12.bin"
model3_path = "model/Model_groups_3class_seed22.bin"

def get_model_path(filename):
    """Get model path from local file only"""
    if os.path.exists(filename):
        return filename
    else:
        raise FileNotFoundError(
            f"Model file '{filename}' not found in current directory.\n"
            f"Please ensure the file exists in: {os.getcwd()}"
        )

# Check if local model exists, raise error if not found
if not os.path.exists(LOCAL_MODEL_DIR) or not os.path.exists(os.path.join(LOCAL_MODEL_DIR, "config.json")):
    raise FileNotFoundError(
        f"Local model not found at: {LOCAL_MODEL_DIR}\n"
        f"Please download the model by running: python download_model.py"
    )
model_path = LOCAL_MODEL_DIR
print(f"Using local model from: {model_path}", flush=True)

# Suppress warnings during model loading
print("Loading tokenizer...", flush=True)
with warnings.catch_warnings():
    warnings.simplefilter("ignore")
    tokenizer = AutoTokenizer.from_pretrained(model_path)
print("Tokenizer loaded.", flush=True)

# Load Model 1 from local path
print("Loading Model 1...", flush=True)
with warnings.catch_warnings():
    warnings.simplefilter("ignore")
    model_1 = AutoModelForSequenceClassification.from_pretrained(model_path, num_labels=41)
print("Loading Model 1 weights...", flush=True)
model_1.load_state_dict(torch.load(model1_path, map_location=device))
model_1.to(device).eval()
# Optimize model for inference (torch.compile available in PyTorch 2.0+)
if torch.cuda.is_available() and hasattr(torch, 'compile'):
    try:
        model_1 = torch.compile(model_1, mode='reduce-overhead')
        print("Model 1 compiled for faster inference.", flush=True)
    except Exception as e:
        print(f"Warning: Could not compile Model 1: {e}", flush=True)
print("Model 1 loaded.", flush=True)

# Load Model 2
print("Loading Model 2...", flush=True)
model2_file = get_model_path(model2_path)
with warnings.catch_warnings():
    warnings.simplefilter("ignore")
    model_2 = AutoModelForSequenceClassification.from_pretrained(model_path, num_labels=41)
print("Loading Model 2 weights...", flush=True)
model_2.load_state_dict(torch.load(model2_file, map_location=device))
model_2.to(device).eval()
# Optimize model for inference (torch.compile available in PyTorch 2.0+)
if torch.cuda.is_available() and hasattr(torch, 'compile'):
    try:
        model_2 = torch.compile(model_2, mode='reduce-overhead')
        print("Model 2 compiled for faster inference.", flush=True)
    except Exception as e:
        print(f"Warning: Could not compile Model 2: {e}", flush=True)
print("Model 2 loaded.", flush=True)

# Load Model 3
print("Loading Model 3...", flush=True)
model3_file = get_model_path(model3_path)
with warnings.catch_warnings():
    warnings.simplefilter("ignore")
    model_3 = AutoModelForSequenceClassification.from_pretrained(model_path, num_labels=41)
print("Loading Model 3 weights...", flush=True)
model_3.load_state_dict(torch.load(model3_file, map_location=device))
model_3.to(device).eval()
# Optimize model for inference (torch.compile available in PyTorch 2.0+)
if torch.cuda.is_available() and hasattr(torch, 'compile'):
    try:
        model_3 = torch.compile(model_3, mode='reduce-overhead')
        print("Model 3 compiled for faster inference.", flush=True)
    except Exception as e:
        print(f"Warning: Could not compile Model 3: {e}", flush=True)
print("Model 3 loaded. All models ready!", flush=True)


# --- Label Mapping and Text Cleaning ---
label_mapping = [
    '13B', '30B', '65B', '7B', 'GLM130B', 'bloom_7b',
    'bloomz', 'cohere', 'davinci', 'dolly', 'dolly-v2-12b',
    'flan_t5_base', 'flan_t5_large', 'flan_t5_small',
    'flan_t5_xl', 'flan_t5_xxl', 'gemma-7b-it', 'gemma2-9b-it',
    'gpt-3.5-turbo', 'gpt-35', 'gpt4', 'gpt4o',
    'gpt_j', 'gpt_neox', 'human', 'llama3-70b', 'llama3-8b',
    'mixtral-8x7b', 'opt_1.3b', 'opt_125m', 'opt_13b',
    'opt_2.7b', 'opt_30b', 'opt_350m', 'opt_6.7b',
    'opt_iml_30b', 'opt_iml_max_1.3b', 't0_11b', 't0_3b',
    'text-davinci-002', 'text-davinci-003'
]

# Default excluded indices (bloomz, davinci, dolly, dolly-v2-12b) - same as backup_app.py
default_excluded = [6, 8, 9, 10]

# label_mapping = [
#     '13B', '30B', '65B', '7B', 'GLM130B', 'bloom_7b',
#     None, 'cohere', None, None, 'dolly-v2-12b',
#     'flan_t5_base', 'flan_t5_large', 'flan_t5_small',
#     'flan_t5_xl', 'flan_t5_xxl', 'gemma-7b-it', 'gemma2-9b-it',
#     'gpt-3.5-turbo', 'gpt-35', 'gpt4', 'gpt4o',
#     'gpt_j', 'gpt_neox', 'human', 'llama3-70b', 'llama3-8b',
#     'mixtral-8x7b', 'opt_1.3b', 'opt_125m', 'opt_13b',
#     'opt_2.7b', 'opt_30b', 'opt_350m', 'opt_6.7b',
#     'opt_iml_30b', 'opt_iml_max_1.3b', 't0_11b', 't0_3b',
#     'text-davinci-002', 'text-davinci-003'
# ]

def extract_text_from_file(file_path: str) -> str:
    """Extract text from various document formats"""
    if not file_path:
        return ""
    
    # Ensure file_path is a string and exists
    if not isinstance(file_path, str):
        file_path = str(file_path)
    
    if not os.path.exists(file_path):
        return f"Error: File not found at path: {file_path}"
    
    file_extension = os.path.splitext(file_path)[1].lower()
    text = ""
    
    try:
        if file_extension == '.pdf':
            try:
                # Use PyMuPDF for text extraction
                doc = pymupdf.open(file_path)
                for page in doc:
                    try:
                        page_text = page.get_text()
                        if page_text:
                            text += page_text + "\n"
                    except Exception as e:
                        # Continue with next page if one fails
                        print(f"Warning: Could not extract text from page {page.number + 1}: {str(e)}", flush=True)
                        continue
                doc.close()
            except Exception as e:
                return f"Error reading PDF file: {str(e)}"
        
        elif file_extension in ['.docx', '.doc']:
            doc = Document(file_path)
            for paragraph in doc.paragraphs:
                text += paragraph.text + "\n"
            # Also extract text from tables
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        text += cell.text + " "
                    text += "\n"
        
        elif file_extension in ['.pptx', '.ppt']:
            prs = Presentation(file_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        
        elif file_extension in ['.xlsx', '.xls']:
            df = pd.read_excel(file_path, sheet_name=None)
            for sheet_name, sheet_df in df.items():
                text += f"\n--- Sheet: {sheet_name} ---\n"
                text += sheet_df.to_string() + "\n"
        
        elif file_extension in ['.csv']:
            df = pd.read_csv(file_path)
            text = df.to_string()
        
        elif file_extension in ['.txt', '.md', '.rtf']:
            with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                text = file.read()
        
        else:
            # Try to read as text file
            try:
                with open(file_path, 'r', encoding='utf-8', errors='ignore') as file:
                    text = file.read()
            except:
                return f"Unsupported file format: {file_extension}. Please upload PDF, DOCX, PPTX, XLSX, CSV, or TXT files."
        
        if not text.strip():
            return f"Error: No text could be extracted from the file. The file might be empty, corrupted, or contain only images."
        
        return text.strip()
    
    except Exception as e:
        return f"Error extracting text from file: {str(e)}"

def clean_text(text: str) -> str:
    text = re.sub(r'\s{2,}', ' ', text)
    text = re.sub(r'\s+([,.;:?!])', r'\1', text)
    return text

def split_into_sentences(text: str) -> list:
    """Split text into sentences using spaCy (multilingual), with NLTK and regex fallbacks"""
    # Try spaCy first (best for multilingual support)
    if SPACY_AVAILABLE and nlp is not None:
        try:
            doc = nlp(text)
            sentences = [sent.text.strip() for sent in doc.sents]
            # Filter out empty sentences
            sentences = [s for s in sentences if s.strip()]
            if sentences:
                return sentences
        except Exception as e:
            print(f"Warning: spaCy sentence splitting failed: {e}, trying NLTK fallback...", flush=True)
    
    # Fallback to NLTK
    try:
        sentences = sent_tokenize(text)
        cleaned = [s.strip() for s in sentences if s.strip()]
        # If NLTK returns only one sentence for a long text, it might have failed
        if len(cleaned) == 1 and len(text) > 200 and ('.' in text or '!' in text or '?' in text):
            raise ValueError("NLTK may have failed to split properly")
        if cleaned:
            return cleaned
    except Exception as e:
        print(f"Warning: NLTK sentence tokenizer failed: {e}, using regex fallback...", flush=True)
    
    # Final fallback to regex-based splitting
    # Improved regex that handles more cases:
    # - Split on . ! ? followed by whitespace (one or more spaces/newlines)
    # - Don't require capital letter after (handles lowercase starts, numbers, etc.)
    # - Handle multiple spaces/newlines/tabs
    sentences = re.split(r'(?<=[.!?])\s+', text)
    # Filter out empty sentences and strip whitespace
    sentences = [s.strip() for s in sentences if s.strip()]
    # If still no sentences found, try more aggressive splitting
    if not sentences or (len(sentences) == 1 and len(text) > 100):
        # Try splitting on period followed by space, newline, or end of string
        sentences = re.split(r'\.(?=\s+|$)', text)
        # Add period back to sentences that were split (except the last one if it already has one)
        result = []
        for i, s in enumerate(sentences):
            s = s.strip()
            if s:
                # Add period if it was split off (not the last sentence or if it doesn't end with punctuation)
                if i < len(sentences) - 1 and not s.endswith(('.', '!', '?')):
                    s += '.'
                result.append(s)
        sentences = result
    return [s for s in sentences if s.strip()]

def generate_per_sentence_debug_info(sentence_data, sentence_ai_percentages, ai_indices, overall_ai_percentage, excluded_indices, sentences):
    """Generate Per-Sentence Debug Information HTML using current calculation (sum of AI probabilities)"""
    debug_info = f"""
    <div style="padding: 20px; background-color: #f5f5f5; border-radius: 8px; max-height: 1750px; overflow-y: auto;">
        <h2 style="margin-top: 0; color: #333;">Per-Sentence Debug Information</h2>
        <div style="margin-bottom: 15px; padding: 10px; background-color: #e3f2fd; border-left: 4px solid #2196F3; border-radius: 4px; font-size: 12px;">
            <strong>Calculation Info:</strong> Overall AI percentage (from per-sentence): {overall_ai_percentage:.2f}% | 
            Total sentences: {len(sentences)} | 
            Sentences marked as AI: {len(ai_indices)} | 
            AI indices: {sorted(ai_indices) if ai_indices else 'None'}
        </div>
    """
    
    # Filter out empty sentences for display (but keep indices aligned)
    non_empty_sentence_data = [(idx, data) for idx, data in enumerate(sentence_data) if idx < len(sentences)]
    
    for display_idx, (orig_idx, (sentence, human_prob_sent, total_ai_prob_sent, ai_percentage)) in enumerate(non_empty_sentence_data, 1):
        if not sentence.strip():
            continue
            
        # Get full probability data for this sentence
        cleaned_sent = clean_text(sentence)
        inputs = tokenizer(cleaned_sent, return_tensors="pt", truncation=True, padding=True, max_length=512).to(device)
        
        with torch.no_grad():
            logits_1 = model_1(**inputs).logits
            logits_2 = model_2(**inputs).logits
            logits_3 = model_3(**inputs).logits
            
            softmax_1 = torch.softmax(logits_1, dim=1)
            softmax_2 = torch.softmax(logits_2, dim=1)
            softmax_3 = torch.softmax(logits_3, dim=1)
            
            averaged_probabilities = (softmax_1 + softmax_2 + softmax_3) / 3
            sent_probs = averaged_probabilities[0]
        
        # Get all probabilities for display
        sent_prob_list = []
        for prob_idx in range(len(sent_probs)):
            if prob_idx < len(label_mapping):
                model_name = label_mapping[prob_idx]
                prob = sent_probs[prob_idx].item()
                sent_prob_list.append((prob_idx, model_name, prob))
        
        sent_prob_list.sort(key=lambda x: x[2], reverse=True)
        
        # Find max indices for display purposes
        sent_probs_for_decision = sent_probs.clone()
        for prob_idx in excluded_indices:
            if 0 <= prob_idx < len(sent_probs_for_decision):
                sent_probs_for_decision[prob_idx] = 0
        sent_max_index = torch.argmax(sent_probs_for_decision).item()
        
        sent_ai_probs_for_max = sent_probs.clone()
        for prob_idx in excluded_indices:
            if 0 <= prob_idx < len(sent_ai_probs_for_max):
                sent_ai_probs_for_max[prob_idx] = 0
        if 24 not in excluded_indices:
            sent_ai_probs_for_max[24] = 0
        sent_max_ai_index = sent_max_index if sent_max_index != 24 else torch.argmax(sent_ai_probs_for_max).item() if sent_ai_probs_for_max.max() > 0 else sent_max_index
        
        # Calculate total probability
        sent_total_prob = sum(p[2] for p in sent_prob_list) if sent_prob_list else 0.0
        
        # Determine sentence type: AI if max_ai_prob > human_prob, Human otherwise
        sent_type = "AI" if orig_idx in ai_indices else "Human"
        sent_color = "#FF5733" if sent_type == "AI" else "#4CAF50"
        
        debug_info += f"""
        <div style="margin-bottom: 20px; padding: 15px; background-color: white; border-left: 4px solid {sent_color}; border-radius: 4px;">
            <div style="display: flex; justify-content: space-between; margin-bottom: 10px;">
                <strong style="color: {sent_color}; font-size: 16px;">Sentence {display_idx} ({sent_type})</strong>
                <span style="font-size: 13px; color: #666;">
                    Human: {human_prob_sent:.4f} | AI: {total_ai_prob_sent:.4f}
                </span>
            </div>
            <div style="margin-bottom: 10px; padding: 10px; background-color: #f9f9f9; border-radius: 4px; font-size: 13px; color: #333;">
                "{html.escape(sentence[:300])}{'...' if len(sentence) > 300 else ''}"
            </div>
            <div style="overflow-x: auto;">
                <table style="width: 100%; border-collapse: collapse; font-family: monospace; font-size: 11px;">
                    <thead>
                        <tr style="background-color: {sent_color}; color: white;">
                            <th style="padding: 6px; text-align: left; border: 1px solid #ddd;">Rank</th>
                            <th style="padding: 6px; text-align: left; border: 1px solid #ddd;">Model</th>
                            <th style="padding: 6px; text-align: right; border: 1px solid #ddd;">Probability</th>
                            <th style="padding: 6px; text-align: right; border: 1px solid #ddd;">Percentage</th>
                            <th style="padding: 6px; text-align: center; border: 1px solid #ddd;">Excluded</th>
                        </tr>
                    </thead>
                    <tbody>
        """
        
        for rank, (idx, model_name, prob) in enumerate(sent_prob_list, 1):
            percentage = (prob / sent_total_prob * 100) if sent_total_prob > 0 else 0.0
            if percentage > 0 and percentage % 1 != 0:
                percentage = math.ceil(percentage)
            else:
                percentage = int(percentage)
            is_excluded = idx in excluded_indices
            excluded_marker = "✓" if is_excluded else ""
            row_color = "#ffebee" if is_excluded else "#ffffff"
            
            debug_info += f"""
                        <tr style="background-color: {row_color};">
                            <td style="padding: 4px; border: 1px solid #ddd;">{rank}</td>
                            <td style="padding: 4px; border: 1px solid #ddd; font-weight: {'bold' if idx == sent_max_index or idx == sent_max_ai_index else 'normal'}">{model_name}</td>
                            <td style="padding: 4px; border: 1px solid #ddd; text-align: right;">{prob:.6f}</td>
                            <td style="padding: 4px; border: 1px solid #ddd; text-align: right;">{percentage}%</td>
                            <td style="padding: 4px; border: 1px solid #ddd; text-align: center; color: {'#d32f2f' if is_excluded else '#666'}">{excluded_marker}</td>
                        </tr>
            """
        
        debug_info += """
                    </tbody>
                </table>
            </div>
        </div>
        """
    
    debug_info += """
        <div style="margin-top: 15px; padding: 10px; background-color: #fff3cd; border-left: 4px solid #ffc107; border-radius: 4px; font-size: 12px;">
            <strong>Note:</strong> Excluded models are marked with ✓ and highlighted in red. Models with bold names are the detected max probability models.
        </div>
    </div>
    """
    
    return debug_info

def classify_sentence(sentence: str, excluded_indices=None):
    """Classify a single sentence and return human and total AI probabilities"""
    if excluded_indices is None:
        excluded_indices = [6, 8, 9, 10]  # Default excluded models (human is included by default)
    
    if not sentence.strip():
        return 0.0, 0.0
    
    cleaned = clean_text(sentence)
    inputs = tokenizer(cleaned, return_tensors="pt", truncation=True, padding=True, max_length=512).to(device)
    
    with torch.no_grad():
        logits_1 = model_1(**inputs).logits
        logits_2 = model_2(**inputs).logits
        logits_3 = model_3(**inputs).logits
        
        softmax_1 = torch.softmax(logits_1, dim=1)
        softmax_2 = torch.softmax(logits_2, dim=1)
        softmax_3 = torch.softmax(logits_3, dim=1)
        
        averaged_probabilities = (softmax_1 + softmax_2 + softmax_3) / 3
        probabilities = averaged_probabilities[0]
    
    # Get human probability (index 24) - set to 0 if human is excluded
    if 24 in excluded_indices:
        human_prob = 0.0
    else:
        human_prob = probabilities[24].item()
    
    # Sum all AI probabilities (excluding human and excluded indices) instead of just max
    # This better reflects the overall AI probability
    ai_probs_sum = probabilities.clone()
    for idx in excluded_indices:
        if 0 <= idx < len(ai_probs_sum):
            ai_probs_sum[idx] = 0
    # Also exclude human (24) from AI sum calculation if not already excluded
    if 24 not in excluded_indices:
        ai_probs_sum[24] = 0
    total_ai_prob = ai_probs_sum.sum().item()
    
    return human_prob, total_ai_prob

def highlight_ai_text(text: str, overall_ai_percentage: float = 50.0, progress=None, excluded_indices=None, return_sentence_data=False):
    """Highlight AI-detected sentences in red and human sentences in green
    Uses percentile-based approach to match overall classification distribution
    Returns: (highlighted_text, ai_char_count, human_char_count) or 
             (highlighted_text, ai_char_count, human_char_count, sentence_data) if return_sentence_data=True"""
    if excluded_indices is None:
        excluded_indices = [6, 8, 9, 10]  # Default excluded models (human is included by default)
    
    sentences = split_into_sentences(text)
    
    # First pass: calculate AI percentages for all sentences
    sentence_ai_percentages = []
    sentence_texts = []
    sentence_data = []  # Store sentence data for debug info
    total_sentences = len(sentences)
    
    for idx, sentence in enumerate(sentences):
        if progress and total_sentences > 0:
            # Update progress: 0.3 to 0.8 (50% of remaining progress)
            progress_val = 0.3 + (idx / total_sentences) * 0.5
            progress(progress_val, desc=f"Analyzing sentences ({idx + 1}/{total_sentences})...")
        human_prob, total_ai_prob = classify_sentence(sentence, excluded_indices)
        total_human_ai_prob = human_prob + total_ai_prob
        if total_human_ai_prob > 0:
            ai_percentage = (total_ai_prob / total_human_ai_prob) * 100
        else:
            ai_percentage = 50.0
        sentence_ai_percentages.append(ai_percentage)
        sentence_texts.append(sentence)
    
        # Store data for debug info if needed
        if return_sentence_data:
            sentence_data.append((sentence, human_prob, total_ai_prob, ai_percentage))
    
    # Calculate threshold based on percentile to match overall distribution
    # If overall is 74% AI, mark top 74% of sentences (by AI%) as AI
    if len(sentence_ai_percentages) > 0:
        # Sort sentences by AI percentage
        sorted_indices = sorted(range(len(sentence_ai_percentages)), 
                               key=lambda i: sentence_ai_percentages[i], 
                               reverse=True)
        # Calculate how many sentences should be marked as AI
        # Use proper rounding: if AI percentage < 50%, use floor (round down) to favor human
        # If AI percentage >= 50%, use ceil (round up) to favor AI
        if overall_ai_percentage < 50.0:
            # When AI is less than 50%, round down to favor human highlighting
            num_ai_sentences = int(len(sentences) * (overall_ai_percentage / 100.0))
        else:
            # When AI is 50% or more, round up to favor AI highlighting
            num_ai_sentences = math.ceil(len(sentences) * (overall_ai_percentage / 100.0))
        # Ensure we don't mark all sentences as AI if overall_ai_percentage < 50%
        if overall_ai_percentage < 50.0 and num_ai_sentences == len(sentences) and len(sentences) > 1:
            num_ai_sentences = len(sentences) - 1
        # Ensure at least one sentence is marked as AI if overall_ai_percentage > 50%
        if overall_ai_percentage > 50.0 and num_ai_sentences == 0 and len(sentences) > 0:
            num_ai_sentences = 1
        ai_indices = set(sorted_indices[:num_ai_sentences])
    else:
        ai_indices = set()
    
    # Second pass: highlight sentences and count characters while preserving original text structure
    if progress:
        progress(0.8, desc="Highlighting text...")
    
    ai_char_count = 0
    human_char_count = 0
    
    # Reconstruct the text by replacing sentences in the original text while preserving structure
    highlighted_text = text
    # Process sentences in order, replacing each occurrence in the original text
    for i, sentence in enumerate(sentence_texts):
        escaped_sentence = html.escape(sentence)
        char_count = len(sentence)
        if i in ai_indices:
            highlighted_sentence = f'<span class="ai-highlight">{escaped_sentence}</span>'
            ai_char_count += char_count
        else:
            highlighted_sentence = f'<span class="human-highlight">{escaped_sentence}</span>'
            human_char_count += char_count
        
        # Replace the sentence in the original text (only first occurrence to handle duplicates correctly)
        escaped_for_regex = re.escape(sentence)
        highlighted_text = re.sub(escaped_for_regex, highlighted_sentence, highlighted_text, count=1)
    
    if return_sentence_data:
        return highlighted_text, ai_char_count, human_char_count, (sentence_data, sentence_ai_percentages, ai_indices, overall_ai_percentage, sentence_texts)
    return highlighted_text, ai_char_count, human_char_count

newline_to_space = Replace(Regex(r'\s*\n\s*'), " ")
join_hyphen_break = Replace(Regex(r'(\w+)[--]\s*\n\s*(\w+)'), r"\1\2")
tokenizer.backend_tokenizer.normalizer = Sequence([
    tokenizer.backend_tokenizer.normalizer,
    join_hyphen_break,
    newline_to_space,
    Strip()
])

# Initialize Language Tool for grammar checking
# Cache for LanguageTool instances per language
language_tools_cache = {}

# Set seed for langdetect to ensure consistent results
DetectorFactory.seed = 0

def detect_language_code_for_grammar(text: str) -> str:
    """Detect the language of the text and return LanguageTool language code
    Returns: LanguageTool language code (e.g., 'en-US', 'de-DE', 'fr-FR')
    Note: This is different from detect_language() which returns True/False for English validation.
    """
    if not text or not text.strip():
        return 'en-US'  # Default to English
    
    try:
        # Detect language using langdetect
        detected_lang = detect(text)
        
        # Map langdetect language codes to LanguageTool language codes
        language_map = {
            'en': 'en-US',  # English -> US English
            'de': 'de-DE',  # German
            'fr': 'fr-FR',  # French
            'es': 'es-ES',  # Spanish
            'it': 'it-IT',  # Italian
            'pt': 'pt-PT',  # Portuguese
            'nl': 'nl-NL',  # Dutch
            'pl': 'pl-PL',  # Polish
            'ru': 'ru-RU',  # Russian
            'zh': 'zh',     # Chinese
            'ja': 'ja',     # Japanese
            'ko': 'ko',     # Korean
            'ar': 'ar',     # Arabic
            'hi': 'hi',     # Hindi
        }
        
        # Return mapped language or default to en-US
        return language_map.get(detected_lang, 'en-US')
    except Exception as e:
        # If detection fails, default to English
        print(f"Language detection failed: {e}, defaulting to en-US", flush=True)
        return 'en-US'

def get_language_tool(language_code='en-US'):
    """Get or create a LanguageTool instance for the specified language"""
    if language_code not in language_tools_cache:
        try:
            language_tools_cache[language_code] = language_tool_python.LanguageTool(language_code)
            print(f"Language Tool initialized for {language_code}.", flush=True)
        except Exception as e:
            print(f"Warning: Could not initialize Language Tool for {language_code}: {e}", flush=True)
            return None
    return language_tools_cache[language_code]

# Initialize default English tool
try:
    tool = get_language_tool('en-US')
    print("Language Tool initialized successfully.", flush=True)
except Exception as e:
    print(f"Warning: Could not initialize Language Tool: {e}", flush=True)
    tool = None

# Initialize NLTK resources
try:
    nltk.download('punkt', quiet=True)
    nltk.download('punkt_tab', quiet=True)
    nltk.download('stopwords', quiet=True)
    nltk.download('averaged_perceptron_tagger', quiet=True)
    print("NLTK resources initialized successfully.", flush=True)
except Exception as e:
    print(f"Warning: Could not download all NLTK resources: {e}", flush=True)
    # Try to download them individually
    try:
        nltk.download('punkt', quiet=True)
    except:
        pass
    try:
        nltk.download('punkt_tab', quiet=True)
    except:
        pass
    try:
        nltk.download('stopwords', quiet=True)
    except:
        pass
    try:
        nltk.download('averaged_perceptron_tagger', quiet=True)
    except:
        pass

# --- Grammar Checking Functions ---
def check_grammar(text: str, language_code=None):
    """Check grammar and return errors with suggestions
    If language_code is None, automatically detects the language from the text."""
    if not text or not text.strip():
        return "", 0, 0
    
    # Auto-detect language if not provided
    if language_code is None:
        language_code = detect_language_code_for_grammar(text)
    
    current_tool = get_language_tool(language_code)
    if current_tool is None:
        return '<div style="color: orange;">Grammar checker not available for this language. Please install language-tool-python.</div>', 0, 0
    
    try:
        matches = current_tool.check(text)
        total_errors = len(matches)
        
        if total_errors == 0:
            result_html = '<div style="color: #4CAF50; font-weight: bold; padding: 15px; border-radius: 8px; background-color: rgba(76, 175, 80, 0.1);">✓ No grammar errors found!</div>'
            return result_html, 0, 0
        
        # Calculate error rate
        word_count = len(word_tokenize(text))
        error_rate = (total_errors / word_count * 100) if word_count > 0 else 0
        
        # Build HTML with original text, underlining only the error spans in red
        result_parts = [f'''
        <div class="grammar-summary">
            Found {total_errors} grammar error(s) (Error rate: {error_rate:.2f}%)
        </div>
        <div class="grammar-text">
        ''']

        # Sort errors by offset for sequential processing
        matches_sorted = sorted(matches, key=lambda x: x.offset)
        
        # Create highlighted (underlined) text by walking through the string once
        last_idx = 0
        for match in matches_sorted:
            start = match.offset
            # Use 'length' attribute (language_tool_python uses 'length' not 'errorLength')
            error_length = getattr(match, 'length', 1)
            end = match.offset + error_length

            # Add normal text before the error
            if start > last_idx:
                normal_segment = text[last_idx:start]
                result_parts.append(html.escape(normal_segment))

            # Add error text with whole word highlighting
            error_text = text[start:end]
            # Expand to whole word
            word_start = start
            while word_start > 0 and text[word_start - 1] not in ' \t\n\r':
                word_start -= 1
            word_end = end
            while word_end < len(text) and text[word_end] not in ' \t\n\r':
                word_end += 1
            
            # Exclude trailing punctuation marks (., ?, !, ", ', :, ;, etc.)
            # Include both regular quotes (" ') and smart/curly quotes (" " ' ')
            # Also include other quote-like characters and check Unicode category
            import unicodedata
            punctuation_chars = '.,?!;:"\'"\u201c\u201d\u2018\u2019\u201a\u201b\u201e\u201f\u2039\u203a()[]{}'
            actual_word_end = word_end
            # Keep trimming punctuation from the end until we hit a non-punctuation character
            while actual_word_end > word_start:
                char = text[actual_word_end - 1]
                # Check if it's in our punctuation list or is a Unicode punctuation/quote category
                char_category = unicodedata.category(char) if char else ''
                is_punctuation = (char in punctuation_chars or 
                                char_category.startswith('P') or  # Any punctuation category
                                char in '\u0022\u0027\u201c\u201d\u2018\u2019\u201a\u201b\u201e\u201f\u2039\u203a\u00ab\u00bb\u301d\u301e\u301f')
                if is_punctuation:
                    actual_word_end -= 1
                else:
                    break
            
            whole_word = text[word_start:actual_word_end]
            result_parts.append(
                f'<span class="grammar-error-word">{html.escape(whole_word)}</span>'
            )

            last_idx = end

        # Add any remaining normal text after the last error
        if last_idx < len(text):
            remaining = text[last_idx:]
            result_parts.append(html.escape(remaining))

        result_parts.append('</div>')  # Close grammar-text
        result_html = ''.join(result_parts)
        
        return result_html, total_errors, error_rate
    except Exception as e:
        return f'<div style="color: red;">Error checking grammar: {html.escape(str(e))}</div>', 0, 0

def add_grammar_highlights_to_text(text: str, highlighted_text: str, language_code=None):
    """Add grammar error highlights to existing highlighted text (with AI/human highlights)
    If language_code is None, automatically detects the language from the text.
    Returns: highlighted_text_with_grammar, grammar_error_count"""
    if not text or not text.strip():
        return highlighted_text, 0
    
    # Auto-detect language if not provided
    if language_code is None:
        language_code = detect_language_code_for_grammar(text)
    
    current_tool = get_language_tool(language_code)
    if current_tool is None:
        return highlighted_text, 0
    
    try:
        matches = current_tool.check(text)
        total_errors = len(matches)
        
        if total_errors == 0:
            return highlighted_text, 0
        
        # Filter out less important errors (whitespace issues, minor formatting)
        important_errors = []
        skip_messages = [
            'whitespace',
            'repeated a whitespace',
            'repeated whitespace',
            'consecutive spaces',
            'too many consecutive spaces',
            'space on both sides',
            'quote symbol',
            'closing parenthesis',
            'opening parenthesis',
            'before the closing',
            'after the opening',
            'punctuation',
            'put a space',
        ]
        
        for match in matches:
            message = str(match.message).lower()
            # Skip whitespace and minor formatting errors
            if any(skip_term in message for skip_term in skip_messages):
                continue
            # Skip very short errors (likely false positives) - but allow single character errors for spelling
            error_length = getattr(match, 'length', 1)
            error_text = text[match.offset:match.offset + error_length] if match.offset + error_length <= len(text) else ''
            # Skip if it's just whitespace
            if error_text.strip() == '':
                continue
            # Skip single character errors that are just punctuation/whitespace
            if error_length == 1 and error_text in ' \t\n':
                continue
            important_errors.append(match)
        
        if len(important_errors) == 0:
            return highlighted_text, 0
        
        # Sort errors by offset (reverse order for insertion to maintain positions)
        matches_sorted = sorted(important_errors, key=lambda x: x.offset, reverse=True)
        
        result_text = highlighted_text
        successfully_inserted_count = 0
        
        # Split text into sentences to find which sentence contains each error
        sentences = split_into_sentences(text)
        
        # Build sentence boundaries
        sentence_boundaries = []
        current_pos = 0
        for sentence in sentences:
            # Find sentence in original text
            sentence_start = text.find(sentence.strip(), current_pos)
            if sentence_start == -1:
                sentence_start = current_pos
            sentence_end = sentence_start + len(sentence)
            sentence_boundaries.append((sentence_start, sentence_end, sentence))
            current_pos = sentence_end
        
        # Process grammar errors in reverse order to maintain correct positions
        for match in matches_sorted:
            start = match.offset
            error_length = getattr(match, 'length', 1)
            end = match.offset + error_length
            
            # Skip if positions are out of range
            if start >= len(text) or end > len(text):
                continue
            
            message = match.message
            replacements = match.replacements[:3] if hasattr(match, 'replacements') and match.replacements else []
            
            # Extract the error text from original text and expand to whole word
            error_text = text[start:end]
            if not error_text.strip():  # Skip if empty after stripping
                continue
            
            # Expand to whole word: find word boundaries (start and end of word)
            # Find the start of the word (go backwards until space or start of text)
            word_start = start
            while word_start > 0 and text[word_start - 1] not in ' \t\n\r':
                word_start -= 1
            
            # Find the end of the word (go forwards until space or end of text)
            word_end = end
            while word_end < len(text) and text[word_end] not in ' \t\n\r':
                word_end += 1
            
            # Exclude trailing punctuation marks (., ?, !, ", ', :, ;, etc.)
            # Trim punctuation from the end of the word
            # Include both regular quotes (" ') and smart/curly quotes (" " ' ')
            # Also include other quote-like characters and check Unicode category
            import unicodedata
            punctuation_chars = '.,?!;:"\'"\u201c\u201d\u2018\u2019\u201a\u201b\u201e\u201f\u2039\u203a()[]{}'
            actual_word_end = word_end
            # Keep trimming punctuation from the end until we hit a non-punctuation character
            while actual_word_end > word_start:
                char = text[actual_word_end - 1]
                # Check if it's in our punctuation list or is a Unicode punctuation/quote category
                char_category = unicodedata.category(char) if char else ''
                is_punctuation = (char in punctuation_chars or 
                                char_category.startswith('P') or  # Any punctuation category
                                char in '\u0022\u0027\u201c\u201d\u2018\u2019\u201a\u201b\u201e\u201f\u2039\u203a\u00ab\u00bb\u301d\u301e\u301f')
                if is_punctuation:
                    actual_word_end -= 1
                else:
                    break
            
            # Use the whole word (excluding trailing punctuation) instead of just the error portion
            whole_word = text[word_start:actual_word_end]
            escaped_error_text = html.escape(whole_word)
            
            # Escape message and replacements for HTML attributes
            msg_escaped = html.escape(str(message)).replace('"', '&quot;')
            repls_escaped = '|'.join([html.escape(str(r)).replace('|', '&#124;') for r in replacements])
            
            # Create grammar error span - mark whole word in red
            grammar_span = f'<span class="grammar-error-word" data-grammar-message="{msg_escaped}" data-grammar-suggestions="{repls_escaped}" style="cursor: pointer;">{escaped_error_text}</span>'
            
            # Find which sentence contains this error
            containing_sentence = None
            sentence_start_pos = None
            for sent_start, sent_end, sentence in sentence_boundaries:
                if sent_start <= start < sent_end:
                    containing_sentence = sentence
                    sentence_start_pos = sent_start
                    break
            
            if containing_sentence:
                # Find this sentence in the highlighted HTML
                escaped_sentence = html.escape(containing_sentence)
                # Look for the sentence within a highlight span
                sentence_pattern = f'<span class="(?:ai|human)-highlight">({re.escape(escaped_sentence)})</span>'
                sentence_match = re.search(sentence_pattern, result_text)
                
                if sentence_match:
                    # Found the sentence in HTML
                    sentence_html_start = sentence_match.start(1)
                    sentence_html_end = sentence_match.end(1)
                    
                    # Calculate position of error within the sentence (in original text)
                    # Use word_start instead of start to account for whole word expansion
                    error_offset_in_sentence = word_start - sentence_start_pos
                    
                    # Get the sentence HTML content
                    sentence_html = result_text[sentence_html_start:sentence_html_end]
                    
                    # Build a character-to-position mapping for the sentence HTML
                    # This maps each character position in plain text to its position in HTML
                    plain_to_html_map = []
                    html_pos = 0
                    plain_pos = 0
                    
                    while html_pos < len(sentence_html):
                        if sentence_html[html_pos] == '<':
                            # Skip HTML tag
                            tag_end = sentence_html.find('>', html_pos)
                            if tag_end == -1:
                                break
                            html_pos = tag_end + 1
                        else:
                            # This HTML position corresponds to this plain text position
                            plain_to_html_map.append((plain_pos, sentence_html_start + html_pos))
                            plain_pos += 1
                            html_pos += 1
                    
                    # Find the HTML positions for the error start and end (use whole_word)
                    error_start_plain = error_offset_in_sentence
                    error_end_plain = error_offset_in_sentence + len(whole_word)
                    
                    # Find the closest HTML positions (handle cases where exact match might not exist)
                    error_html_start_pos = None
                    error_html_end_pos = None
                    
                    # Find start position (closest match)
                    for plain_pos, html_pos in plain_to_html_map:
                        if plain_pos >= error_start_plain:
                            error_html_start_pos = html_pos
                            break
                    
                    # If not found, use the first position
                    if error_html_start_pos is None and len(plain_to_html_map) > 0:
                        error_html_start_pos = plain_to_html_map[0][1]
                    
                    # Find end position (closest match) - use whole_word length
                    error_end_plain = error_offset_in_sentence + len(whole_word)
                    for plain_pos, html_pos in plain_to_html_map:
                        if plain_pos >= error_end_plain:
                            error_html_end_pos = html_pos
                            break
                    
                    # If not found, calculate from start + whole word length
                    if error_html_end_pos is None:
                        if error_html_start_pos is not None:
                            # Count characters from start to find end
                            chars_counted = 0
                            search_pos = error_html_start_pos - sentence_html_start
                            while search_pos < len(sentence_html) and chars_counted < len(whole_word):
                                if sentence_html[search_pos] == '<':
                                    tag_end = sentence_html.find('>', search_pos)
                                    if tag_end == -1:
                                        break
                                    search_pos = tag_end + 1
                                else:
                                    chars_counted += 1
                                    error_html_end_pos = sentence_html_start + search_pos + 1
                                    search_pos += 1
                        else:
                            # Last resort: use sentence end
                            error_html_end_pos = sentence_html_end
                    
                    # If we found the positions, extract and replace
                    if error_html_start_pos is not None and error_html_end_pos is not None and error_html_end_pos > error_html_start_pos:
                        # Check if this position is already inside a grammar error span
                        before_text = result_text[:error_html_start_pos]
                        if 'grammar-error-word' in before_text or 'grammar-error-underline' in before_text:
                            # Check if we're inside an unclosed span
                            last_grammar_open = before_text.rfind('<span class="grammar-error-word"')
                            if last_grammar_open == -1:
                                last_grammar_open = before_text.rfind('<span class="grammar-error-underline"')
                            last_grammar_close = before_text.rfind('</span>', last_grammar_open)
                            if last_grammar_open > last_grammar_close:
                                # We're inside a grammar span, skip this error
                                continue
                        
                        # Extract the error text from HTML to verify
                        html_error_segment = result_text[error_html_start_pos:error_html_end_pos]
                        plain_error_segment = re.sub(r'<[^>]+>', '', html_error_segment)
                        
                        # Verify we have the right text (allow some flexibility for whitespace)
                        # Use whole_word instead of error_text for matching
                        if (whole_word.strip() in plain_error_segment.strip() or 
                            plain_error_segment.strip() in whole_word.strip() or
                            whole_word in plain_error_segment or
                            escaped_error_text in html_error_segment):
                            # Replace the segment with grammar span - this ensures the entire word is wrapped
                            result_text = result_text[:error_html_start_pos] + grammar_span + result_text[error_html_end_pos:]
                            successfully_inserted_count += 1
                        else:
                            # Fallback: try to find the whole word in the sentence at the correct position
                            plain_sentence = re.sub(r'<[^>]+>', '', sentence_html)
                            if whole_word in plain_sentence:
                                # Find the occurrence closest to our target position
                                plain_start = plain_sentence.find(whole_word, max(0, error_offset_in_sentence - 5))
                                if plain_start == -1:
                                    plain_start = plain_sentence.find(whole_word)
                                
                                if plain_start != -1 and abs(plain_start - error_offset_in_sentence) < 50:
                                    # Map to HTML position
                                    if plain_start < len(plain_to_html_map):
                                        html_start = plain_to_html_map[plain_start][1]
                                        # Find end position
                                        html_end = html_start
                                        chars_counted = 0
                                        search_pos = html_start - sentence_html_start
                                        while search_pos < len(sentence_html) and chars_counted < len(whole_word):
                                            if sentence_html[search_pos] == '<':
                                                tag_end = sentence_html.find('>', search_pos)
                                                if tag_end == -1:
                                                    break
                                                search_pos = tag_end + 1
                                            else:
                                                chars_counted += 1
                                                html_end = sentence_html_start + search_pos + 1
                                                search_pos += 1
                                        
                                        if html_end > html_start:
                                            result_text = result_text[:html_start] + grammar_span + result_text[html_end:]
                                            successfully_inserted_count += 1
                                    else:
                                        # Last resort: simple replacement
                                        if escaped_error_text in sentence_html:
                                            sentence_with_grammar = sentence_html.replace(escaped_error_text, grammar_span, 1)
                                            result_text = result_text[:sentence_html_start] + sentence_with_grammar + result_text[sentence_html_end:]
                                            successfully_inserted_count += 1
                                else:
                                    # Last resort: simple replacement
                                    if escaped_error_text in sentence_html:
                                        sentence_with_grammar = sentence_html.replace(escaped_error_text, grammar_span, 1)
                                        result_text = result_text[:sentence_html_start] + sentence_with_grammar + result_text[sentence_html_end:]
                    else:
                        # Fallback: simple text replacement when position mapping fails
                        plain_sentence = re.sub(r'<[^>]+>', '', sentence_html)
                        if error_text in plain_sentence:
                            # Find the occurrence closest to our target position (use whole_word)
                            plain_start = plain_sentence.find(whole_word, max(0, error_offset_in_sentence - 5))
                            if plain_start == -1:
                                plain_start = plain_sentence.find(whole_word)
                            
                            if plain_start != -1 and abs(plain_start - error_offset_in_sentence) < 50:
                                # Map to HTML using the mapping we built
                                if plain_start < len(plain_to_html_map):
                                    html_start = plain_to_html_map[plain_start][1]
                                    # Find end position (use whole_word length)
                                    if plain_start + len(whole_word) < len(plain_to_html_map):
                                        html_end = plain_to_html_map[plain_start + len(whole_word)][1]
                                    else:
                                        # Calculate end by counting
                                        html_end = html_start
                                        chars_counted = 0
                                        search_pos = html_start - sentence_html_start
                                        while search_pos < len(sentence_html) and chars_counted < len(whole_word):
                                            if sentence_html[search_pos] == '<':
                                                tag_end = sentence_html.find('>', search_pos)
                                                if tag_end == -1:
                                                    break
                                                search_pos = tag_end + 1
                                            else:
                                                chars_counted += 1
                                                html_end = sentence_html_start + search_pos + 1
                                                search_pos += 1
                                    
                                    if html_end > html_start:
                                        result_text = result_text[:html_start] + grammar_span + result_text[html_end:]
                                        successfully_inserted_count += 1
                            else:
                                # Last resort: replace first occurrence
                                if escaped_error_text in sentence_html:
                                    sentence_with_grammar = sentence_html.replace(escaped_error_text, grammar_span, 1)
                                    result_text = result_text[:sentence_html_start] + sentence_with_grammar + result_text[sentence_html_end:]
            else:
                # Fallback: try to find the whole word directly in HTML (not in a tag)
                pattern = f'(?<=>)([^<]*?){re.escape(escaped_error_text)}([^<]*?)(?=<)'
                match_obj = re.search(pattern, result_text)
                if match_obj:
                    before_text = match_obj.group(1)
                    after_text = match_obj.group(2)
                    replacement = before_text + grammar_span + after_text
                    result_text = result_text[:match_obj.start()] + replacement + result_text[match_obj.end():]
                    successfully_inserted_count += 1
        
        return result_text, successfully_inserted_count
        
    except Exception as e:
        # If grammar checking fails, return original highlighted text
        return highlighted_text, 0

# --- TOEFL Rubric Scoring Functions ---
def calculate_toefl_score(text: str):
    """Calculate TOEFL writing rubric scores based on official TOEFL iBT rubric"""
    if not text or not text.strip():
        return {}
    
    sentences = sent_tokenize(text)
    words = word_tokenize(text.lower())
    word_count = len([w for w in words if w.isalnum()])
    
    # Remove stopwords for analysis
    try:
        stop_words = set(stopwords.words('english'))
        content_words = [w for w in words if w.isalnum() and w not in stop_words]
    except:
        content_words = [w for w in words if w.isalnum()]
    
    # TOEFL iBT Writing Rubric Criteria (0-5 scale)
    
    # 1. Development (0-5 points)
    # Based on: relevant and adequately elaborated explanations, details, examples
    development_score = 2.5  # Base score
    
    # Length factor (TOEFL typically expects 150-300 words for good development)
    if word_count >= 200 and word_count <= 400:
        development_score += 1.0
    elif word_count >= 150 and word_count < 200:
        development_score += 0.5
    elif word_count >= 400:
        development_score += 1.5
    elif word_count < 100:
        development_score -= 1.0
    
    # Sentence complexity and elaboration
    avg_sentence_length = word_count / len(sentences) if sentences else 0
    if 15 <= avg_sentence_length <= 30:
        development_score += 0.5
    elif avg_sentence_length < 10:
        development_score -= 0.5
    
    # Check for examples and details (words like "for example", "such as", "specifically")
    example_indicators = ['example', 'instance', 'specifically', 'particularly', 'namely', 'such as', 'including']
    example_count = sum(1 for word in words if word in example_indicators)
    if example_count >= 2:
        development_score += 0.5
    
    development_score = min(5.0, max(0.0, development_score))
    
    # 2. Organization (0-5 points)
    # Based on: clear structure, transitions, coherence
    organization_score = 2.5  # Base score
    
    transition_words = ['however', 'therefore', 'furthermore', 'moreover', 'additionally', 
                       'consequently', 'thus', 'hence', 'nevertheless', 'although', 'because',
                       'first', 'second', 'third', 'finally', 'in conclusion', 'for example',
                       'in addition', 'on the other hand', 'meanwhile', 'subsequently']
    transition_count = sum(1 for word in words if word in transition_words)
    
    if len(sentences) > 0:
        transition_ratio = transition_count / len(sentences)
        if transition_ratio >= 0.3:
            organization_score += 1.5
        elif transition_ratio >= 0.2:
            organization_score += 1.0
        elif transition_ratio >= 0.1:
            organization_score += 0.5
        else:
            organization_score -= 0.5
    
    # Check for clear structure (introductory phrases, conclusion indicators)
    structure_indicators = ['first', 'second', 'third', 'finally', 'in conclusion', 'to summarize', 'in summary']
    structure_count = sum(1 for word in words if word in structure_indicators)
    if structure_count >= 2:
        organization_score += 0.5
    
    organization_score = min(5.0, max(0.0, organization_score))
    
    # 3. Language Use (0-5 points)
    # Based on: grammar, vocabulary, sentence variety
    grammar_score = 5.0
    grammar_errors = 0
    current_tool = get_language_tool('en-US')  # TOEFL is English-specific
    if current_tool:
        try:
            matches = current_tool.check(text)
            grammar_errors = len(matches)
            error_rate = (grammar_errors / word_count) if word_count > 0 else 0
            # More lenient scoring - TOEFL allows some errors
            if error_rate < 0.02:  # Less than 2% error rate
                grammar_score = 5.0
            elif error_rate < 0.05:  # Less than 5% error rate
                grammar_score = 4.0
            elif error_rate < 0.10:  # Less than 10% error rate
                grammar_score = 3.0
            else:
                grammar_score = max(0, 5 - (error_rate * 30))
        except:
            pass
    
    # Vocabulary diversity and sophistication
    unique_words = len(set(content_words))
    vocab_diversity = (unique_words / len(content_words)) if content_words else 0
    
    # Check for academic/sophisticated vocabulary
    academic_words = ['analyze', 'demonstrate', 'illustrate', 'examine', 'evaluate', 'assess',
                     'significant', 'substantial', 'considerable', 'essential', 'crucial', 'vital']
    academic_count = sum(1 for word in words if word in academic_words)
    
    vocab_score = 2.5  # Base
    if vocab_diversity > 0.65:
        vocab_score += 1.5
    elif vocab_diversity > 0.50:
        vocab_score += 1.0
    elif vocab_diversity > 0.40:
        vocab_score += 0.5
    else:
        vocab_score -= 0.5
    
    if academic_count >= 3:
        vocab_score += 0.5
    
    vocab_score = min(5.0, max(0.0, vocab_score))
    
    # Sentence variety
    sentence_variety_score = 3.0  # Base
    if 12 <= avg_sentence_length <= 25:
        sentence_variety_score += 1.0
    elif avg_sentence_length > 25:
        sentence_variety_score += 0.5
    elif avg_sentence_length < 8:
        sentence_variety_score -= 1.0
    
    # Check for sentence structure variety (simple, compound, complex)
    complex_indicators = ['although', 'because', 'while', 'whereas', 'despite', 'although', 'if', 'when']
    complex_count = sum(1 for word in words if word in complex_indicators)
    if complex_count >= 2:
        sentence_variety_score += 0.5
    
    sentence_variety_score = min(5.0, max(0.0, sentence_variety_score))
    
    language_score = (grammar_score * 0.4 + vocab_score * 0.35 + sentence_variety_score * 0.25)
    
    # Overall score (average of three categories, rounded to nearest 0.5 for TOEFL scale)
    overall_score = (development_score + organization_score + language_score) / 3
    # Round to nearest 0.5 (TOEFL scores are typically in 0.5 increments)
    overall_score = round(overall_score * 2) / 2
    
    scores = {
        'development': round(development_score, 1),
        'organization': round(organization_score, 1),
        'language': round(language_score, 1),
        'overall': round(overall_score, 1),
        'word_count': word_count,
        'sentence_count': len(sentences),
        'grammar_errors': grammar_errors if tool else 0,
        'transition_count': transition_count,
        'vocab_diversity': round(vocab_diversity * 100, 1)
    }
    
    return scores

def generate_toefl_feedback(text: str, scores: dict):
    """Generate detailed TOEFL feedback based on scores, matching official TOEFL feedback style"""
    if not scores or not text:
        return "", ""
    
    overall = scores.get('overall', 0)
    dev_score = scores.get('development', 0)
    org_score = scores.get('organization', 0)
    lang_score = scores.get('language', 0)
    word_count = scores.get('word_count', 0)
    grammar_errors = scores.get('grammar_errors', 0)
    transition_count = scores.get('transition_count', 0)
    
    sentences = sent_tokenize(text)
    words = word_tokenize(text.lower())
    
    # Analyze text characteristics for more specific feedback
    try:
        stop_words = set(stopwords.words('english'))
        content_words = [w for w in words if w.isalnum() and w not in stop_words]
        unique_words = len(set(content_words))
        vocab_diversity = (unique_words / len(content_words)) if content_words else 0
    except:
        vocab_diversity = 0.5
    
    # Generate overall feedback (more detailed and contextual)
    strengths = []
    improvements = []
    
    # Development analysis
    if dev_score >= 4.0:
        strengths.append("clearly explained ideas")
        if word_count >= 200:
            strengths.append("adequate elaboration")
    elif dev_score >= 3.0:
        if word_count < 150:
            improvements.append("expand your ideas with more details and examples")
        else:
            improvements.append("add more specific examples to support your points")
    else:
        improvements.append("develop your ideas more thoroughly with detailed explanations and concrete examples")
    
    # Organization analysis
    if org_score >= 4.0:
        strengths.append("well-organized structure")
        if transition_count >= len(sentences) * 0.2:
            strengths.append("effective use of transition words")
    elif org_score >= 3.0:
        if transition_count < len(sentences) * 0.1:
            improvements.append("add more transition words to improve flow between ideas")
        else:
            improvements.append("strengthen organization with clearer structure")
    else:
        improvements.append("improve organization with better structure and more connecting phrases")
    
    # Language use analysis
    if lang_score >= 4.0:
        strengths.append("varied sentence structures")
        if vocab_diversity > 0.5:
            strengths.append("appropriate vocabulary")
        if grammar_errors <= 2:
            strengths.append("accurate grammar")
    elif lang_score >= 3.0:
        if grammar_errors > 5:
            improvements.append("review grammar to reduce errors")
        if vocab_diversity < 0.4:
            improvements.append("use more varied vocabulary")
        else:
            improvements.append("enhance language use with more precise word choice")
    else:
        improvements.append("improve language use with better grammar, more varied vocabulary, and sentence variety")
    
    # Build overall feedback (matching the style from the image)
    if strengths and overall >= 3.5:
        overall_feedback = f"You {', '.join(strengths[:3])}"
        if len(strengths) > 3:
            overall_feedback += f", and {strengths[3]}"
        
        if improvements:
            overall_feedback += f". To strengthen your essay, {', '.join(improvements[:2])}"
            if len(improvements) > 2:
                overall_feedback += f", especially when {improvements[2].replace('add more', 'adding more').replace('improve', 'improving').replace('enhance', 'enhancing')}"
            overall_feedback += ", and ensure precise language and clear connections between sections."
        else:
            overall_feedback += "."
        
        overall_feedback += " Keep building on these strengths, and with deeper elaboration your writing will reach an even higher level."
    elif strengths:
        overall_feedback = f"Your essay shows {', '.join(strengths[:2])}"
        if improvements:
            overall_feedback += f". To improve, focus on {', '.join(improvements[:2])}"
        overall_feedback += "."
    else:
        overall_feedback = f"Your essay addresses the topic but needs improvement. Focus on {', '.join(improvements[:3])} to strengthen your writing."
    
    # Generate task-specific feedback (TOEFL Academic Discussion Task style)
    if overall >= 4.0:
        task_feedback = "The essay is a generally successful response as it provides relevant and adequately elaborated explanations"
        if dev_score >= 4.0:
            task_feedback += " with good development of ideas"
        if org_score >= 4.0:
            task_feedback += " and clear organization"
        if lang_score >= 4.0:
            task_feedback += ". The student uses a variety of syntactic structures and appropriate word choice"
        if grammar_errors <= 2:
            task_feedback += ", with few lexical or grammatical errors"
        if improvements:
            task_feedback += f". However, some sections could benefit from {improvements[0].replace('add more', 'more').replace('improve', 'improved').replace('enhance', 'enhanced')} and deeper elaboration."
        else:
            task_feedback += "."
    elif overall >= 3.0:
        task_feedback = "The essay provides relevant explanations"
        if dev_score < 3.5:
            task_feedback += ", though some sections need more elaboration and specific examples"
        if org_score < 3.5:
            task_feedback += ". Organization could be improved with better structure and transitions"
        if lang_score < 3.5:
            task_feedback += ", and language use shows some errors that affect clarity"
        task_feedback += ". With more practice and attention to detail, the response can reach a higher level."
    else:
        task_feedback = "The essay addresses the topic but needs significant improvement in development, organization, and language use"
        if improvements:
            task_feedback += f". Focus on {improvements[0]} and {improvements[1] if len(improvements) > 1 else 'overall clarity'}"
        task_feedback += " to strengthen your writing."
    
    return overall_feedback, task_feedback

def format_toefl_rubric(scores: dict, text: str = ""):
    """Format TOEFL rubric scores as HTML matching the official TOEFL design"""
    if not scores:
        return ""
    
    dev_score = scores.get('development', 0)
    org_score = scores.get('organization', 0)
    lang_score = scores.get('language', 0)
    overall = scores.get('overall', 0)
    
    # Generate feedback
    overall_feedback, task_feedback = generate_toefl_feedback(text, scores) if text else ("", "")
    
    # Determine score color
    if overall >= 4.0:
        score_color = "#4CAF50"  # Green
    elif overall >= 3.0:
        score_color = "#8BC34A"  # Light green
    elif overall >= 2.0:
        score_color = "#FFC107"  # Yellow
    else:
        score_color = "#FF9800"  # Orange
    
    html_output = f'''
    <div style="padding: 0; font-family: -apple-system, BlinkMacSystemFont, 'Segoe UI', Roboto, sans-serif;">
        <!-- TOEFL Score Header -->
        <div style="display: flex; align-items: center; gap: 20px; margin-bottom: 30px; padding: 20px; background: #f8f9fa; border-radius: 12px;">
            <!-- Circular Badge -->
            <div style="width: 100px; height: 100px; border-radius: 50%; background: {score_color}; border: 4px solid white; box-shadow: 0 4px 12px rgba(0,0,0,0.15); display: flex; align-items: center; justify-content: center; flex-shrink: 0;">
                <span style="font-size: 36px; font-weight: bold; color: white;">{overall:.1f}</span>
            </div>
            
            <!-- Score Info -->
            <div style="flex: 1;">
                <h1 style="margin: 0; font-size: 32px; font-weight: 700; color: #1a1a1a;">TOEFL Score</h1>
                <p style="margin: 8px 0; font-size: 24px; font-weight: 600; color: #666;">{overall:.1f} / 5</p>
                <p style="margin: 8px 0 0 0; font-size: 14px; color: #666;">
                    Score calculated using <a href="https://www.ets.org/toefl/test-takers/ibt/about/content/writing" target="_blank" style="color: #0066cc; text-decoration: underline;">TOEFL rubric</a>
                </p>
            </div>
        </div>
        
        <!-- Overall Feedback Section -->
        <div style="background: #f5f5f5; border-radius: 12px; padding: 20px; margin-bottom: 20px; border: 1px solid #e0e0e0;">
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 15px;">
                <h2 style="margin: 0; font-size: 20px; font-weight: 600; color: #1a1a1a;">Overall Feedback</h2>
                <span style="font-size: 18px; color: #666;">▲</span>
            </div>
            <p style="margin: 0; font-size: 15px; line-height: 1.6; color: #333;">{html.escape(overall_feedback) if overall_feedback else "No feedback available."}</p>
        </div>
        
        <!-- TOEFL Academic Discussion Task Section -->
        <div style="background: #f5f5f5; border-radius: 12px; padding: 20px; border: 1px solid #e0e0e0;">
            <div style="display: flex; justify-content: space-between; align-items: center; margin-bottom: 15px;">
                <div style="display: flex; align-items: center; gap: 10px;">
                    <h2 style="margin: 0; font-size: 20px; font-weight: 600; color: #1a1a1a;">
                        <a href="https://www.ets.org/toefl/test-takers/ibt/about/content/writing" target="_blank" style="color: #0066cc; text-decoration: underline;">TOEFL Academic Discussion Task</a>
                    </h2>
                    <span style="font-size: 16px; font-weight: 600; color: #666;">{overall:.1f} / 5</span>
                </div>
                <span style="font-size: 18px; color: #666;">▲</span>
            </div>
            <p style="margin: 0; font-size: 15px; line-height: 1.6; color: #333;">{html.escape(task_feedback) if task_feedback else "No task-specific feedback available."}</p>
        </div>
        
        <!-- Detailed Breakdown (Collapsible) -->
        <div style="margin-top: 20px; background: white; border-radius: 12px; padding: 20px; border: 1px solid #e0e0e0;">
            <h3 style="margin-top: 0; font-size: 18px; font-weight: 600; color: #1a1a1a; margin-bottom: 20px;">Detailed Breakdown</h3>
            
            <div style="margin-bottom: 20px;">
                <div style="display: flex; justify-content: space-between; margin-bottom: 8px;">
                    <span style="font-weight: 600; color: #333;">Development</span>
                    <span style="font-weight: 600; color: #666;">{dev_score:.1f} / 5.0</span>
                </div>
                <div style="background: #e0e0e0; height: 8px; border-radius: 4px; overflow: hidden;">
                    <div style="background: {score_color}; height: 100%; width: {dev_score/5*100}%; transition: width 0.3s;"></div>
                </div>
                <p style="margin: 8px 0 0 0; font-size: 13px; color: #666;">Measures how well you develop your ideas with details and examples.</p>
            </div>
            
            <div style="margin-bottom: 20px;">
                <div style="display: flex; justify-content: space-between; margin-bottom: 8px;">
                    <span style="font-weight: 600; color: #333;">Organization</span>
                    <span style="font-weight: 600; color: #666;">{org_score:.1f} / 5.0</span>
                </div>
                <div style="background: #e0e0e0; height: 8px; border-radius: 4px; overflow: hidden;">
                    <div style="background: {score_color}; height: 100%; width: {org_score/5*100}%; transition: width 0.3s;"></div>
                </div>
                <p style="margin: 8px 0 0 0; font-size: 13px; color: #666;">Measures how well you organize your writing with clear structure and transitions.</p>
            </div>
            
            <div style="margin-bottom: 10px;">
                <div style="display: flex; justify-content: space-between; margin-bottom: 8px;">
                    <span style="font-weight: 600; color: #333;">Language Use</span>
                    <span style="font-weight: 600; color: #666;">{lang_score:.1f} / 5.0</span>
                </div>
                <div style="background: #e0e0e0; height: 8px; border-radius: 4px; overflow: hidden;">
                    <div style="background: {score_color}; height: 100%; width: {lang_score/5*100}%; transition: width 0.3s;"></div>
                </div>
                <p style="margin: 8px 0 0 0; font-size: 13px; color: #666;">Measures grammar, vocabulary, and sentence variety.</p>
            </div>
        </div>
    </div>
    '''
    
    return html_output

# --- AI Writing Feedback Functions ---
# Local LLM model for writing feedback
_local_llm_model = None
_local_llm_tokenizer = None
_local_llm_device = device
_current_model_name = None

def load_local_llm(model_name: str = "microsoft/Phi-3-mini-4k-instruct"):
    """Load local LLM model for text generation"""
    global _local_llm_model, _local_llm_tokenizer, _current_model_name
    
    # If model is already loaded with the same name, return it
    if _local_llm_model is not None and _current_model_name == model_name:
        return _local_llm_model, _local_llm_tokenizer
    
    # If a different model is requested, we'd need to reload (for now, just use existing)
    # In production, you might want to allow model switching
    if _local_llm_model is not None and _current_model_name != model_name:
        print(f"Warning: Model {_current_model_name} already loaded. To use {model_name}, restart the application.", flush=True)
        return _local_llm_model, _local_llm_tokenizer
    
    try:
        from transformers import AutoModelForCausalLM, AutoTokenizer, pipeline
        from transformers.utils import logging as transformers_logging
        
        # Suppress transformers warnings about flash-attention
        transformers_logging.set_verbosity_error()
        
        print(f"Loading local LLM model: {model_name}...", flush=True)
        
        # Try to load from local path first, then from HuggingFace
        local_model_path = os.path.join("model", "llm")
        if os.path.exists(local_model_path) and os.path.exists(os.path.join(local_model_path, "config.json")):
            model_path = local_model_path
            print(f"Using local model from: {model_path}", flush=True)
        else:
            # Use HuggingFace Hub for faster downloads with resume capability
            try:
                from huggingface_hub import snapshot_download
                
                # Get cache directory from environment or use default
                cache_dir = os.environ.get("HF_HOME", os.path.expanduser("~/.cache/huggingface"))
                
                # Download to local path with resume capability
                print(f"Downloading model from HuggingFace (this may take a while, download will resume if interrupted)...", flush=True)
                print(f"Model will be cached at: {local_model_path}", flush=True)
                
                # Create local model directory
                os.makedirs(local_model_path, exist_ok=True)
                
                # Download with resume capability
                downloaded_path = snapshot_download(
                    repo_id=model_name,
                    local_dir=local_model_path,
                    local_dir_use_symlinks=False,  # Use actual files, not symlinks
                    resume_download=True,  # Resume interrupted downloads
                    cache_dir=cache_dir
                )
                model_path = local_model_path
                print(f"Model downloaded successfully to: {model_path}", flush=True)
            except ImportError:
                # Fallback to transformers download if huggingface_hub not available
                print(f"huggingface_hub not available, using standard download (slower, no resume)...", flush=True)
                model_path = model_name
            except Exception as e:
                print(f"Error downloading with huggingface_hub: {e}, falling back to standard download...", flush=True)
                model_path = model_name
        
        with warnings.catch_warnings():
            warnings.simplefilter("ignore")
            # Suppress all warnings including UserWarnings about flash-attention
            warnings.filterwarnings("ignore", message=".*flash-attention.*")
            warnings.filterwarnings("ignore", message=".*flash_attn.*")
            
            _local_llm_tokenizer = AutoTokenizer.from_pretrained(model_path, trust_remote_code=True)
            _local_llm_model = AutoModelForCausalLM.from_pretrained(
                model_path,
                torch_dtype=torch.float16 if torch.cuda.is_available() else torch.float32,
                device_map="auto" if torch.cuda.is_available() else None,
                trust_remote_code=True,
                low_cpu_mem_usage=True,
                attn_implementation="eager"  # Use eager attention to avoid flash-attention warning
            )
            
            if not torch.cuda.is_available():
                _local_llm_model = _local_llm_model.to(_local_llm_device)
            
            _local_llm_model.eval()
            
            # Set pad token if not set
            if _local_llm_tokenizer.pad_token is None:
                _local_llm_tokenizer.pad_token = _local_llm_tokenizer.eos_token
        
        _current_model_name = model_name
        print(f"Local LLM model loaded successfully on {_local_llm_device}", flush=True)
        return _local_llm_model, _local_llm_tokenizer
        
    except Exception as e:
        print(f"Error loading local LLM model: {str(e)}", flush=True)
        print("Falling back to rule-based feedback only.", flush=True)
        return None, None

def generate_ai_feedback(text: str, model_name: str = None):
    """Generate AI-powered writing feedback using local LLM"""
    if not text or not text.strip():
        return None
    
    try:
        global _local_llm_model, _local_llm_tokenizer
        
        # Load model if not already loaded
        if _local_llm_model is None or _local_llm_tokenizer is None:
            model_name = model_name or "microsoft/Phi-3-mini-4k-instruct"
            load_local_llm(model_name)
        
        if _local_llm_model is None or _local_llm_tokenizer is None:
            return None
        
        # Truncate text if too long (keep it reasonable for the model)
        max_input_length = 2000
        if len(text) > max_input_length:
            text = text[:max_input_length] + "..."
        
        # Detect model type and format prompt accordingly
        model_name_lower = (model_name or "microsoft/Phi-3-mini-4k-instruct").lower()
        
        if "phi" in model_name_lower or "phi-3" in model_name_lower:
            # Phi-3 format
            prompt = f"""<|system|>
You are an expert writing tutor providing detailed, constructive feedback on student writing. Analyze the text and provide comprehensive feedback covering:

1. Content & Ideas: Evaluate clarity, depth, and development of ideas
2. Organization: Assess structure, flow, and logical progression
3. Language & Style: Review grammar, vocabulary, sentence variety, and tone
4. Strengths: Identify what the writer does well
5. Areas for Improvement: Provide specific, actionable suggestions

Provide feedback in a clear, encouraging, and constructive manner. Format your response with clear sections.<|end|>
<|user|>
Analyze the following text and provide detailed writing feedback:

{text}<|end|>
<|assistant|>
"""
        elif "llama" in model_name_lower or "mistral" in model_name_lower:
            # Llama/Mistral format
            prompt = f"""<s>[INST] <<SYS>>
You are an expert writing tutor providing detailed, constructive feedback on student writing. Analyze the text and provide comprehensive feedback covering:

1. Content & Ideas: Evaluate clarity, depth, and development of ideas
2. Organization: Assess structure, flow, and logical progression
3. Language & Style: Review grammar, vocabulary, sentence variety, and tone
4. Strengths: Identify what the writer does well
5. Areas for Improvement: Provide specific, actionable suggestions

Provide feedback in a clear, encouraging, and constructive manner. Format your response with clear sections.
<</SYS>>

Analyze the following text and provide detailed writing feedback:

{text} [/INST]"""
        else:
            # Generic format (works for most models)
            prompt = f"""You are an expert writing tutor providing detailed, constructive feedback on student writing. 

Analyze the following text and provide comprehensive feedback covering:
1. Content & Ideas: Evaluate clarity, depth, and development of ideas
2. Organization: Assess structure, flow, and logical progression
3. Language & Style: Review grammar, vocabulary, sentence variety, and tone
4. Strengths: Identify what the writer does well
5. Areas for Improvement: Provide specific, actionable suggestions

Text to analyze:
{text}

Provide detailed writing feedback in a clear, encouraging, and constructive manner:"""
        
        # Tokenize input
        inputs = _local_llm_tokenizer(prompt, return_tensors="pt", truncation=True, max_length=2048).to(_local_llm_device)
        
        # Generate response
        with torch.no_grad():
            outputs = _local_llm_model.generate(
                **inputs,
                max_new_tokens=500,
                temperature=0.7,
                do_sample=True,
                top_p=0.9,
                pad_token_id=_local_llm_tokenizer.eos_token_id,
                eos_token_id=_local_llm_tokenizer.eos_token_id
            )
        
        # Decode response
        generated_text = _local_llm_tokenizer.decode(outputs[0], skip_special_tokens=True)
        
        # Extract only the assistant's response (after the prompt)
        if "<|assistant|>" in generated_text:
            ai_feedback = generated_text.split("<|assistant|>")[-1].strip()
        elif "[/INST]" in generated_text:
            ai_feedback = generated_text.split("[/INST]")[-1].strip()
        else:
            # Fallback: remove the prompt part
            ai_feedback = generated_text[len(prompt):].strip()
        
        # Clean up the feedback - remove special tokens
        ai_feedback = ai_feedback.replace("<|end|>", "").replace("<|endoftext|>", "").strip()
        # Remove any remaining prompt artifacts
        if ai_feedback.startswith("Analyze") and ":" in ai_feedback:
            # Try to find where actual feedback starts
            lines = ai_feedback.split("\n")
            for i, line in enumerate(lines):
                if any(keyword in line.lower() for keyword in ["content", "organization", "language", "strength", "improvement"]):
                    ai_feedback = "\n".join(lines[i:]).strip()
                    break
        
        if not ai_feedback or len(ai_feedback) < 50:
            return None
        
        return ai_feedback
        
    except Exception as e:
        print(f"Error generating AI feedback with local LLM: {str(e)}", flush=True)
        import traceback
        traceback.print_exc()
        return None
    finally:
        cleanup_memory()

def generate_writing_feedback(text: str):
    """Generate comprehensive detailed writing feedback"""
    if not text or not text.strip():
        return ""
    
    sentences = sent_tokenize(text)
    words = word_tokenize(text.lower())
    word_count = len([w for w in words if w.isalnum()])
    char_count = len(text)
    paragraph_count = len([p for p in text.split('\n\n') if p.strip()])
    
    # Detailed analysis data
    analysis = {}
    
    # 1. Length Analysis
    analysis['length'] = {
        'word_count': word_count,
        'char_count': char_count,
        'sentence_count': len(sentences),
        'paragraph_count': paragraph_count,
        'avg_words_per_sentence': round(word_count / len(sentences), 1) if sentences else 0,
        'status': 'good' if 150 <= word_count <= 500 else ('short' if word_count < 150 else 'long')
    }
    
    # 2. Grammar Analysis
    # Use the same filtering logic as add_grammar_highlights_to_text to get accurate count
    grammar_errors = []
    grammar_error_count = 0
    current_tool = get_language_tool('en-US')
    if current_tool:
        try:
            matches = current_tool.check(text)
            # Filter out less important errors (same logic as add_grammar_highlights_to_text)
            skip_messages = [
                'whitespace',
                'repeated a whitespace',
                'repeated whitespace',
                'consecutive spaces',
                'too many consecutive spaces',
                'space on both sides',
                'quote symbol',
                'closing parenthesis',
                'opening parenthesis',
                'before the closing',
                'after the opening',
                'punctuation',
                'put a space',
            ]
            
            important_errors = []
            for match in matches:
                message = str(match.message).lower()
                # Skip whitespace and minor formatting errors
                if any(skip_term in message for skip_term in skip_messages):
                    continue
                # Skip very short errors (likely false positives) - but allow single character errors for spelling
                error_length = getattr(match, 'length', 1)
                error_text = text[match.offset:match.offset + error_length] if match.offset + error_length <= len(text) else ''
                # Skip if it's just whitespace
                if error_text.strip() == '':
                    continue
                # Skip single character errors that are just punctuation/whitespace
                if error_length == 1 and error_text in ' \t\n':
                    continue
                important_errors.append(match)
            
            grammar_error_count = len(important_errors)
            grammar_errors = [(m.offset, m.length, m.message, m.replacements[:3] if hasattr(m, 'replacements') and m.replacements else []) for m in important_errors[:10]]  # Top 10 errors
        except:
            pass
    
    analysis['grammar'] = {
        'error_count': grammar_error_count,
        'error_rate': round((grammar_error_count / word_count * 100), 2) if word_count > 0 else 0,
        'errors': grammar_errors,
        'status': 'excellent' if grammar_error_count == 0 else ('good' if grammar_error_count < 5 else 'needs_improvement')
    }
    
    # 3. Vocabulary Analysis
    try:
        stop_words = set(stopwords.words('english'))
        content_words = [w for w in words if w.isalnum() and w not in stop_words]
        unique_words = len(set(content_words))
        total_content_words = len(content_words)
        diversity = (unique_words / total_content_words) if total_content_words > 0 else 0
        
        # Find repeated words
        word_freq = {}
        for word in content_words:
            word_freq[word] = word_freq.get(word, 0) + 1
        most_repeated = sorted(word_freq.items(), key=lambda x: x[1], reverse=True)[:5]
        
        analysis['vocabulary'] = {
            'unique_words': unique_words,
            'total_content_words': total_content_words,
            'diversity_score': round(diversity * 100, 1),
            'most_repeated': most_repeated,
            'status': 'excellent' if diversity > 0.6 else ('good' if diversity > 0.4 else 'needs_improvement')
        }
    except:
        analysis['vocabulary'] = {'status': 'unknown'}
    
    # 4. Sentence Structure Analysis
    sentence_lengths = [len(word_tokenize(s)) for s in sentences]
    avg_sentence_length = sum(sentence_lengths) / len(sentence_lengths) if sentence_lengths else 0
    min_sentence_length = min(sentence_lengths) if sentence_lengths else 0
    max_sentence_length = max(sentence_lengths) if sentence_lengths else 0
    
    # Count sentence types
    simple_sentences = sum(1 for s in sentences if ',' not in s and ';' not in s and ':' not in s)
    complex_sentences = len(sentences) - simple_sentences
    
    analysis['sentence_structure'] = {
        'avg_length': round(avg_sentence_length, 1),
        'min_length': min_sentence_length,
        'max_length': max_sentence_length,
        'simple_count': simple_sentences,
        'complex_count': complex_sentences,
        'variety_score': round((1 - (max_sentence_length - min_sentence_length) / max_sentence_length) * 100, 1) if max_sentence_length > 0 else 0,
        'status': 'good' if 15 <= avg_sentence_length <= 25 else ('short' if avg_sentence_length < 10 else 'long')
    }
    
    # 5. Organization Analysis
    transition_words = ['however', 'therefore', 'furthermore', 'moreover', 'additionally', 
                       'consequently', 'thus', 'hence', 'nevertheless', 'although', 'because',
                       'first', 'second', 'third', 'finally', 'in conclusion', 'for example',
                       'specifically', 'namely', 'that is', 'in other words', 'to illustrate',
                       'meanwhile', 'subsequently', 'accordingly', 'as a result', 'on the other hand']
    transition_count = sum(1 for word in words if word in transition_words)
    transition_ratio = transition_count / len(sentences) if sentences else 0
    
    # Check for paragraph structure
    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
    avg_paragraph_length = sum(len(word_tokenize(p)) for p in paragraphs) / len(paragraphs) if paragraphs else 0
    
    analysis['organization'] = {
        'transition_count': transition_count,
        'transition_ratio': round(transition_ratio, 2),
        'paragraph_count': paragraph_count,
        'avg_paragraph_length': round(avg_paragraph_length, 1),
        'status': 'good' if transition_ratio >= 0.3 else ('fair' if transition_ratio >= 0.15 else 'needs_improvement')
    }
    
    # 6. Readability Analysis
    # Simple readability metrics
    avg_chars_per_word = char_count / word_count if word_count > 0 else 0
    syllables_estimate = sum(max(1, len([c for c in word if c in 'aeiou'])) for word in words[:100])  # Sample
    avg_syllables = syllables_estimate / min(100, len(words)) if words else 0
    
    # Flesch Reading Ease approximation (simplified)
    asl = avg_sentence_length
    asw = avg_syllables
    flesch_score = 206.835 - (1.015 * asl) - (84.6 * asw)
    
    analysis['readability'] = {
        'flesch_score': round(flesch_score, 1),
        'level': 'very_easy' if flesch_score >= 80 else ('easy' if flesch_score >= 70 else ('fairly_easy' if flesch_score >= 60 else ('standard' if flesch_score >= 50 else ('fairly_difficult' if flesch_score >= 30 else 'difficult')))),
        'avg_chars_per_word': round(avg_chars_per_word, 1)
    }
    
    # Format detailed feedback as HTML
    html_parts = ['<div style="padding: 20px; font-family: Arial, sans-serif;">']
    html_parts.append('<h2 style="color: #667eea; margin-bottom: 25px; border-bottom: 3px solid #667eea; padding-bottom: 10px;">📝 Comprehensive Writing Feedback</h2>')
    
    # Overall Summary
    html_parts.append('<div style="background: linear-gradient(135deg, #667eea 0%, #764ba2 100%); color: white; padding: 20px; border-radius: 10px; margin-bottom: 25px;">')
    html_parts.append('<h3 style="margin-top: 0; color: white;">📊 Overall Statistics</h3>')
    html_parts.append(f'<div style="display: grid; grid-template-columns: repeat(auto-fit, minmax(150px, 1fr)); gap: 15px; margin-top: 15px;">')
    html_parts.append(f'<div><strong>Words:</strong> {word_count}</div>')
    html_parts.append(f'<div><strong>Sentences:</strong> {len(sentences)}</div>')
    html_parts.append(f'<div><strong>Paragraphs:</strong> {paragraph_count}</div>')
    html_parts.append(f'<div><strong>Characters:</strong> {char_count:,}</div>')
    html_parts.append('</div></div>')
    
    # 1. Length Feedback
    len_data = analysis['length']
    status_icon = "✓" if len_data['status'] == 'good' else ("⚠️" if len_data['status'] == 'short' else "ℹ️")
    status_color = "#4CAF50" if len_data['status'] == 'good' else "#FFC107"
    
    html_parts.append(f'''
    <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid {status_color}; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
        <h3 style="margin-top: 0; color: {status_color}; font-size: 18px;">{status_icon} Length & Structure</h3>
        <div style="color: #555; line-height: 1.8;">
            <p><strong>Word Count:</strong> {len_data['word_count']} words</p>
            <p><strong>Average Words per Sentence:</strong> {len_data['avg_words_per_sentence']} words</p>
            <p><strong>Paragraph Count:</strong> {len_data['paragraph_count']} paragraph(s)</p>
            <p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;">
                {html.escape("Your text has " + str(len_data['word_count']) + " words. " + 
                ("This is an appropriate length for most writing tasks." if len_data['status'] == 'good' else 
                ("Your text is quite short. Consider expanding your ideas with more details, examples, or explanations to reach at least 150-300 words for a complete response." if len_data['status'] == 'short' else 
                "Your text is comprehensive and well-developed. Make sure all points are necessary and well-supported.")))}
            </p>
        </div>
        </div>
        ''')
    
    # 2. Grammar Feedback
    gram_data = analysis['grammar']
    status_icon = "✓" if gram_data['status'] == 'excellent' else ("⚠️" if gram_data['status'] == 'good' else "❌")
    status_color = "#4CAF50" if gram_data['status'] == 'excellent' else ("#FFC107" if gram_data['status'] == 'good' else "#FF5733")
    
    html_parts.append(f'''
    <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid {status_color}; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
        <h3 style="margin-top: 0; color: {status_color}; font-size: 18px;">{status_icon} Grammar & Mechanics</h3>
        <div style="color: #555; line-height: 1.8;">
            <p><strong>Grammar Errors Found:</strong> {gram_data['error_count']}</p>
            <p><strong>Error Rate:</strong> {gram_data['error_rate']}%</p>
            {('<p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;"><strong>Recommendation:</strong> ' + html.escape("Excellent! Your writing demonstrates strong grammatical control." if gram_data['status'] == 'excellent' else ("You have a few minor errors. Review the grammar suggestions and make corrections." if gram_data['status'] == 'good' else "There are several grammar errors that need attention. Focus on reviewing basic grammar rules and proofreading carefully.")) + '</p>') if gram_data['error_count'] > 0 else ''}
        </div>
        </div>
        ''')
    
    # 3. Vocabulary Feedback
    if 'diversity_score' in analysis['vocabulary']:
        vocab_data = analysis['vocabulary']
        status_icon = "✓" if vocab_data['status'] == 'excellent' else ("⚠️" if vocab_data['status'] == 'good' else "❌")
        status_color = "#4CAF50" if vocab_data['status'] == 'excellent' else ("#FFC107" if vocab_data['status'] == 'good' else "#FF5733")
        
        repeated_words_text = ", ".join([f'"{word}" ({count}x)' for word, count in vocab_data['most_repeated'][:3]])
        
        html_parts.append(f'''
        <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid {status_color}; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
            <h3 style="margin-top: 0; color: {status_color}; font-size: 18px;">{status_icon} Vocabulary & Word Choice</h3>
            <div style="color: #555; line-height: 1.8;">
                <p><strong>Unique Words:</strong> {vocab_data['unique_words']} out of {vocab_data['total_content_words']} content words</p>
                <p><strong>Vocabulary Diversity Score:</strong> {vocab_data['diversity_score']}%</p>
                <p><strong>Most Repeated Words:</strong> {repeated_words_text}</p>
                <p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;">
                    <strong>Analysis:</strong> {html.escape("Your vocabulary is diverse and engaging. You're using a good variety of words to express your ideas." if vocab_data['status'] == 'excellent' else ("Your vocabulary is adequate but could be more varied. Try using synonyms and more descriptive words to make your writing more engaging." if vocab_data['status'] == 'good' else "Your vocabulary needs more variety. Consider using a thesaurus to find synonyms for frequently repeated words, and try to incorporate more descriptive and precise language."))}
                </p>
            </div>
        </div>
        ''')
    
    # 4. Sentence Structure Feedback
    sent_data = analysis['sentence_structure']
    status_icon = "✓" if sent_data['status'] == 'good' else "⚠️"
    status_color = "#4CAF50" if sent_data['status'] == 'good' else "#FFC107"
    
    html_parts.append(f'''
    <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid {status_color}; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
        <h3 style="margin-top: 0; color: {status_color}; font-size: 18px;">{status_icon} Sentence Structure & Variety</h3>
        <div style="color: #555; line-height: 1.8;">
            <p><strong>Average Sentence Length:</strong> {sent_data['avg_length']} words</p>
            <p><strong>Sentence Length Range:</strong> {sent_data['min_length']} - {sent_data['max_length']} words</p>
            <p><strong>Simple Sentences:</strong> {sent_data['simple_count']} | <strong>Complex Sentences:</strong> {sent_data['complex_count']}</p>
            <p><strong>Variety Score:</strong> {sent_data['variety_score']}%</p>
            <p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;">
                <strong>Recommendation:</strong> {html.escape("Your sentence structure is well-balanced with good variety in length and complexity." if sent_data['status'] == 'good' else ("Your sentences are quite short. Try combining related ideas using conjunctions (and, but, or) or subordinating conjunctions (because, although, while) to create more sophisticated sentence structures." if sent_data['status'] == 'short' else "Some of your sentences are very long and may be difficult to follow. Consider breaking them into shorter, clearer sentences. Aim for 15-25 words per sentence on average."))}
            </p>
        </div>
    </div>
    ''')
    
    # 5. Organization Feedback
    org_data = analysis['organization']
    status_icon = "✓" if org_data['status'] == 'good' else ("⚠️" if org_data['status'] == 'fair' else "❌")
    status_color = "#4CAF50" if org_data['status'] == 'good' else ("#FFC107" if org_data['status'] == 'fair' else "#FF5733")
    
    html_parts.append(f'''
    <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid {status_color}; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
        <h3 style="margin-top: 0; color: {status_color}; font-size: 18px;">{status_icon} Organization & Flow</h3>
        <div style="color: #555; line-height: 1.8;">
            <p><strong>Transition Words Used:</strong> {org_data['transition_count']} ({org_data['transition_ratio']} per sentence)</p>
            <p><strong>Average Paragraph Length:</strong> {org_data['avg_paragraph_length']} words</p>
            <p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;">
                <strong>Analysis:</strong> {html.escape("Your writing is well-organized with good use of transition words to connect ideas smoothly." if org_data['status'] == 'good' else ("Your organization is adequate, but adding more transition words would improve the flow between ideas. Consider using words like 'furthermore', 'however', 'therefore', 'consequently', and 'for example'." if org_data['status'] == 'fair' else "Your writing needs more transition words and phrases to connect ideas better. Use words like 'first', 'second', 'furthermore', 'however', 'therefore', 'in conclusion', and 'for example' to guide readers through your argument."))}
            </p>
        </div>
    </div>
    ''')
    
    # 6. Readability Feedback
    read_data = analysis['readability']
    level_names = {
        'very_easy': 'Very Easy (5th grade)',
        'easy': 'Easy (6th grade)',
        'fairly_easy': 'Fairly Easy (7th-8th grade)',
        'standard': 'Standard (9th-10th grade)',
        'fairly_difficult': 'Fairly Difficult (11th-12th grade)',
        'difficult': 'Difficult (College level)'
    }
    
    html_parts.append(f'''
    <div style="margin-bottom: 20px; padding: 20px; border-left: 5px solid #667eea; background-color: #f8f9fa; border-radius: 8px; box-shadow: 0 2px 4px rgba(0,0,0,0.1);">
        <h3 style="margin-top: 0; color: #667eea; font-size: 18px;">📖 Readability Analysis</h3>
        <div style="color: #555; line-height: 1.8;">
            <p><strong>Flesch Reading Ease Score:</strong> {read_data['flesch_score']}</p>
            <p><strong>Reading Level:</strong> {level_names.get(read_data['level'], read_data['level'])}</p>
            <p><strong>Average Characters per Word:</strong> {read_data['avg_chars_per_word']}</p>
            <p style="margin-top: 10px; padding: 10px; background-color: white; border-radius: 5px;">
                <strong>Interpretation:</strong> Your text is written at a {level_names.get(read_data['level'], read_data['level'])} reading level. 
                {'This is appropriate for most general audiences.' if read_data['level'] in ['fairly_easy', 'standard'] else 
                ('Consider simplifying some sentences for better accessibility.' if read_data['level'] in ['fairly_difficult', 'difficult'] else 
                'You may want to use more sophisticated vocabulary and sentence structures for academic writing.')}
            </p>
        </div>
    </div>
    ''')
    
    # Action Items
    html_parts.append('<div style="margin-top: 30px; padding: 20px; background: linear-gradient(135deg, #f093fb 0%, #f5576c 100%); color: white; border-radius: 10px;">')
    html_parts.append('<h3 style="margin-top: 0; color: white;">🎯 Key Recommendations</h3>')
    html_parts.append('<ul style="line-height: 2; padding-left: 20px;">')
    
    recommendations = []
    if analysis['length']['status'] == 'short':
        recommendations.append("Expand your ideas with more details, examples, or explanations")
    if analysis['grammar']['error_count'] > 0:
        recommendations.append(f"Review and correct {analysis['grammar']['error_count']} grammar error(s)")
    if analysis['vocabulary'].get('status') == 'needs_improvement':
        recommendations.append("Use a thesaurus to find synonyms and increase vocabulary variety")
    if analysis['sentence_structure']['status'] != 'good':
        recommendations.append("Vary your sentence length and structure for better flow")
    if analysis['organization']['status'] != 'good':
        recommendations.append("Add more transition words to improve connections between ideas")
    
    if not recommendations:
        recommendations.append("Continue maintaining the quality of your writing!")
    
    for rec in recommendations:
        html_parts.append(f'<li>{html.escape(rec)}</li>')
    
    html_parts.append('</ul></div>')
    html_parts.append('</div>')
    
    return ''.join(html_parts)

# --- Plagiarism Detection Functions ---
# PrepostSEO API Configuration
PLAGIARISM_API_KEY = os.getenv("PLAGIARISM_API_KEY", "843a7f0eed6d2cbaf0559fe4e151cb328e7096b4")
PLAGIARISM_API_URL = os.getenv("PLAGIARISM_API_URL", "https://www.prepostseo.com/apis/checkPlag")
# Development mode settings
DEV_MODE = os.getenv("DEV_MODE", "false").lower() == "true"
DEBUG = os.getenv("DEBUG", "false").lower() == "true"

def check_plagiarism_api(text: str, exclude_urls: list = None):
    """Check plagiarism using PrepostSEO API
    Maximum words limit: 5,000 words via API"""
    if not text or not text.strip():
        return None
    
    try:
        # Check word count - API has a maximum limit of 5,000 words
        words = text.split()
        word_count = len(words)
        
        # Truncate text if it exceeds 5,000 words
        if word_count > 5000:
            print(f"Warning: Text exceeds 5,000 word limit ({word_count} words). Truncating to first 5,000 words.", flush=True)
            truncated_text = ' '.join(words[:5000])
            text = truncated_text
            word_count = 5000
        
        # Use correct API format: application/x-www-form-urlencoded with key and data parameters
        # This matches the curl format: --data-urlencode 'key=...' --data-urlencode 'data=...'
        headers = {
            "Content-Type": "application/x-www-form-urlencoded"
        }
        
        # Prepare data with key and data parameters (requests.post will automatically URL-encode these)
        # This is equivalent to curl's --data-urlencode
        data = {
            "key": PLAGIARISM_API_KEY,
            "data": text
        }
        
        # Add exclude URLs if provided (if API supports it)
        if exclude_urls and len(exclude_urls) > 0:
            for i, url in enumerate(exclude_urls):
                if url and url.strip():
                    data[f"exclude_urls[{i}]"] = url.strip()
        
        # Make POST request to the correct endpoint
        # Using data=dict automatically URL-encodes values (equivalent to curl --data-urlencode)
        # Increased timeout to 120 seconds to handle longer processing times
        response = requests.post(PLAGIARISM_API_URL, headers=headers, data=data, timeout=120)
        
        if response.status_code == 200:
            try:
                result = response.json()
                return result
            except json.JSONDecodeError:
                print(f"API returned non-JSON response: {response.text[:200]}", flush=True)
                return None
        elif response.status_code == 401:
            # Handle authentication/authorization error
            error_msg = "Plagiarism API authentication failed. Please check your API key."
            print(f"API Error 401 (Unauthorized): {error_msg}", flush=True)
            return {"error": error_msg, "status_code": 401}
        elif response.status_code == 429:
            # Handle rate limit error
            error_msg = "Plagiarism API rate limit exceeded. Please try again later."
            print(f"API Error 429 (Rate Limit): {error_msg}", flush=True)
            return {"error": error_msg, "status_code": 429}
        else:
            # Handle all other non-200 status codes
            try:
                error_response = response.json()
                error_msg = error_response.get('error', error_response.get('message', f'Plagiarism API error (Status {response.status_code})'))
            except (json.JSONDecodeError, AttributeError):
                error_msg = f"Plagiarism API error (Status {response.status_code}). Please try again later."
            print(f"API Error: Status {response.status_code}, Response: {response.text[:500]}", flush=True)
            return {"error": error_msg, "status_code": response.status_code}
    except requests.exceptions.Timeout:
        print("API request timed out", flush=True)
        return None
    except requests.exceptions.RequestException as e:
        print(f"API request error: {str(e)}", flush=True)
        return None
    except Exception as e:
        print(f"Error calling plagiarism API: {str(e)}", flush=True)
        return None

def check_plagiarism(text: str, exclude_urls: list = None):
    """Check for potential plagiarism using PrepostSEO API"""
    if not text or not text.strip():
        return "", 0.0
    
    # Check word count and warn if approaching limit
    words = text.split()
    word_count = len(words)
    warning_msg = ""
    if word_count > 5000:
        # Text will be truncated in the API call, but show a warning in the result
        warning_msg = f'<div style="color: #FFC107; padding: 10px; margin-bottom: 10px; background-color: rgba(255, 193, 7, 0.1); border-radius: 8px; border-left: 4px solid #FFC107;"><strong>⚠️ Warning:</strong> Your text contains {word_count} words. Large texts may be truncated by the detection service.</div>'
    
    # Try to use PrepostSEO API
    api_result = check_plagiarism_api(text, exclude_urls)
    
    # Check if API returned an error (e.g., 401, 429, or any non-200 status)
    if api_result and isinstance(api_result, dict) and api_result.get('error'):
        error_msg = api_result.get('error', 'Plagiarism API error')
        status_code = api_result.get('status_code', 500)
        
        # Determine error color and icon based on status code
        if status_code == 401:
            error_color = "#FF5733"
            error_icon = "🔐"
            error_title = "Authentication Error"
        elif status_code == 429:
            error_color = "#FF9800"
            error_icon = "⏱️"
            error_title = "Rate Limit Exceeded"
        elif status_code in [500, 502, 503]:
            error_color = "#9C27B0"
            error_icon = "🔧"
            error_title = "Server Error"
        elif status_code == 400:
            error_color = "#F44336"
            error_icon = "❌"
            error_title = "Invalid Request"
        else:
            error_color = "#FF5733"
            error_icon = "⚠️"
            error_title = "Plagiarism Check Error"
        
        error_html = f'''
        <div style="padding: 20px; border-radius: 10px; border-left: 4px solid {error_color}; background-color: rgba(255, 87, 51, 0.1);">
            <h2 style="color: {error_color}; margin-top: 0;">{error_icon} {error_title}</h2>
            <p style="color: #333; margin: 10px 0;">{html.escape(error_msg)}</p>
            {f'<p style="color: #666; font-size: 12px; margin-top: 10px;">Status Code: {status_code}</p>' if status_code else ''}
        </div>
        '''
        # Return error result - the calling code should handle this
        return (error_html, 0.0, {"error": error_msg, "status_code": status_code})
    
    if api_result:
        try:
            # Parse API response - new format uses plagPercent and uniquePercent
            plagiarism_percent = float(api_result.get('plagPercent', api_result.get('plagiarism_percentage', 0)))
            unique_percent = float(api_result.get('uniquePercent', api_result.get('unique_percentage', 100)))
            
            # Get sources from the new format
            sources = api_result.get('sources', [])
            # Convert sources to matched_sources format for compatibility
            matched_sources = []
            for source in sources:
                if isinstance(source, dict):
                    matched_sources.append({
                        'url': source.get('link', ''),
                        'similarity': source.get('percent', 0),
                        'count': source.get('count', 0)
                    })
            
            # Also check details array for additional information
            details = api_result.get('details', [])
            if details and not matched_sources:
                # Extract sources from details if available
                for detail in details:
                    if isinstance(detail, dict) and not detail.get('unique', True):
                        display = detail.get('display', {})
                        if isinstance(display, dict):
                            url = display.get('url', '')
                            if url:
                                matched_sources.append({
                                    'url': url,
                                    'similarity': detail.get('percentage', 0),
                                    'count': 1
                                })
            
            # Store full API result for comprehensive report
            api_result['matched_sources'] = matched_sources
            
            # Build result HTML
            if plagiarism_percent > 80:
                color = "#FF5733"
                icon = "❌"
                level = "High"
            elif plagiarism_percent > 60:
                color = "#FFC107"
                icon = "⚠️"
                level = "Moderate"
            elif plagiarism_percent > 30:
                color = "#FFC107"
                icon = "⚠️"
                level = "Low-Moderate"
            else:
                color = "#4CAF50"
                icon = "✓"
                level = "Low"
            
            result_html = warning_msg + f'''
            <div style="padding: 20px; border-radius: 10px; border-left: 4px solid {color}; background-color: rgba(0,0,0,0.02);">
                <h2 style="color: {color}; margin-top: 0;">{icon} Plagiarism Check Results</h2>
                
                <div style="display: flex; gap: 20px; margin: 20px 0;">
                    <div style="flex: 1; padding: 15px; background: rgba(255, 87, 51, 0.1); border-radius: 8px; text-align: center;">
                        <h3 style="margin: 0; color: #FF5733;">Plagiarized</h3>
                        <p style="font-size: 32px; font-weight: bold; margin: 10px 0; color: #FF5733;">{plagiarism_percent:.1f}%</p>
                    </div>
                    <div style="flex: 1; padding: 15px; background: rgba(76, 175, 80, 0.1); border-radius: 8px; text-align: center;">
                        <h3 style="margin: 0; color: #4CAF50;">Unique</h3>
                        <p style="font-size: 32px; font-weight: bold; margin: 10px 0; color: #4CAF50;">{unique_percent:.1f}%</p>
                    </div>
                </div>
                
                <div style="margin: 20px 0;">
                    <div style="background: #e0e0e0; height: 30px; border-radius: 15px; overflow: hidden; display: flex;">
                        <div style="width: {plagiarism_percent}%; background-color: #FF5733; transition: width 0.3s ease;"></div>
                        <div style="width: {unique_percent}%; background-color: #4CAF50; transition: width 0.3s ease;"></div>
                    </div>
                </div>
            '''
            
            if matched_sources and len(matched_sources) > 0:
                result_html += f'''
                <div style="margin-top: 20px;">
                    <h3 style="color: {color};">Matched Sources ({len(matched_sources)})</h3>
                    <ul style="max-height: 300px; overflow-y: auto;">
                '''
                for source in matched_sources[:10]:  # Show first 10 sources
                    url = source.get('url', '')
                    similarity = source.get('similarity', 0)
                    result_html += f'''
                    <li style="margin: 10px 0; padding: 10px; background: rgba(0,0,0,0.02); border-radius: 4px;">
                        <strong>Similarity: {similarity:.1f}%</strong><br>
                        <a href="{html.escape(url)}" target="_blank" style="color: #667eea; word-break: break-all;">{html.escape(url)}</a>
                    </li>
                    '''
                if len(matched_sources) > 10:
                    result_html += f'<li style="color: #666;">... and {len(matched_sources) - 10} more sources</li>'
                result_html += '</ul></div>'
            
            # Add recommendations
            if plagiarism_percent > 80:
                result_html += '''
                <div style="margin-top: 20px; padding: 15px; background: rgba(255, 87, 51, 0.1); border-radius: 8px;">
                    <strong>⚠️ Warning:</strong> Your text shows very high similarity to existing sources. This may indicate plagiarism. 
                    Please ensure proper citation and original writing.
                </div>
                '''
            elif plagiarism_percent > 60:
                result_html += '''
                <div style="margin-top: 20px; padding: 15px; background: rgba(255, 193, 7, 0.1); border-radius: 8px;">
                    <strong>⚠️ Caution:</strong> Your text shows moderate similarity to existing sources. 
                    Ensure you're properly paraphrasing and citing sources.
                </div>
                '''
            elif plagiarism_percent > 30:
                result_html += '''
                <div style="margin-top: 20px; padding: 15px; background: rgba(255, 193, 7, 0.1); border-radius: 8px;">
                    <strong>ℹ️ Note:</strong> Some similarity detected. Review matched sources and ensure proper attribution where needed.
                </div>
                '''
            else:
                result_html += '''
                <div style="margin-top: 20px; padding: 15px; background: rgba(76, 175, 80, 0.1); border-radius: 8px;">
                    <strong>✓ Good:</strong> Your text appears to be mostly original with low similarity to existing sources.
                </div>
                '''
            
            result_html += '</div>'
            
            # Return result_html, plagiarism_percent, and full API result for comprehensive report
            return result_html, plagiarism_percent, api_result
        except Exception as e:
            print(f"Error parsing API response: {str(e)}", flush=True)
            # Fall through to fallback method
    
    # Fallback message if API fails (self-plagiarism detection disabled)
    return '<div style="color: #FFC107; padding: 15px; border-radius: 8px; background-color: rgba(255, 193, 7, 0.1);">⚠️ Plagiarism API is currently unavailable. Please try again later or check your API connection. Self-plagiarism detection is disabled.</div>', 0.0


# --- Citation Detection Functions ---

def detect_citations(text: str):
    """
    Detect common academic citation patterns in the text and return
    an HTML report plus a short Markdown summary string.
    This is a heuristic detector (regex-based), not a full citation parser.
    """
    if not text or not text.strip():
        msg_html = (
            '<div style="color: #666; padding: 10px; background-color: #f5f5f5; '
            'border-radius: 8px;">No text provided for citation analysis.</div>'
        )
        summary = "**Citations Detected:** 0"
        return msg_html, summary

    # Non-URL citation patterns (processed on original text)
    patterns = [
        # APA-style: (Smith, 2020) or (Smith & Jones, 2020, p. 10)
        (r"\([A-Z][A-Za-z\-]+(?:\s+[A-Z][A-Za-z\-]+)*(?:\s*&\s*[A-Z][A-Za-z\-]+(?:\s+[A-Z][A-Za-z\-]+)*)?,\s*\d{4}[a-z]?(?:,[^)]+)?\)", "APA-style"),
        # IEEE / numeric style: [1], [2-4], [10]
        (r"\[\d{1,3}(?:\s*[-,]\s*\d{1,3})*\]", "Numeric-style"),
        # Author-year without parentheses e.g. Smith (2020)
        (r"\b[A-Z][A-Za-z\-]+(?:\s+[A-Z][A-Za-z\-]+)*\s*\(\d{4}[a-z]?\)", "Author-year"),
        # DOI-style references
        (r"\bdoi:\s*10\.\d{4,9}/\S+\b", "DOI"),
    ]

    # Separate URL pattern: we will apply this on a URL-normalized copy of the text
    url_pattern = re.compile(r"https?://\S+")

    # Helper: check if a URL returns HTTP 200 (follows redirects)
    def _url_is_ok(url: str) -> bool:
        try:
            # Use HEAD first to avoid downloading full content
            # allow_redirects=True will follow all redirects (301, 302, 307, 308, etc.)
            # requests library automatically follows up to 30 redirects by default
            resp = requests.head(url, allow_redirects=True, timeout=10)
            # Accept any 2xx status code (200, 201, 202, etc.) as success
            if 200 <= resp.status_code < 300:
                return True
            # Some servers don't support HEAD correctly; fall back to GET
            if resp.status_code in (405, 403, 500):
                resp = requests.get(url, allow_redirects=True, timeout=10)
                # Accept any 2xx status code as success
                return 200 <= resp.status_code < 300
            return False
        except Exception:
            return False

    matches = []
    seen_pairs = set()  # track (label, match) to avoid duplicate rows

    # First, handle non-URL citation patterns on the original text
    for pattern, label in patterns:
        for m in re.finditer(pattern, text):
            span = m.group(0)
            cleaned_match = re.sub(r"\s+", " ", span).strip()
            key = (label, cleaned_match)
            if key in seen_pairs:
                continue
            seen_pairs.add(key)
            matches.append(
                {
                    "label": label,
                    "match": cleaned_match,
                    "start": m.start(),
                    "end": m.end(),
                    "status_ok": None,
                }
            )

    # Next, detect URLs on a cleaned copy of the text.
    # Step 1: collapse newlines to spaces
    normalized_text = text.replace("\n", " ")
    # Step 2: remove spaces immediately before URL structural characters
    # e.g. 'example .com' -> 'example.com', '? via=ihub' (space before '?') -> '?via=ihub'
    normalized_text = re.sub(r"\s+(?=[./?#=&])", "", normalized_text)
    # Step 3: fix common broken 'org' domain split such as 'oclc.or g/science' -> 'oclc.org/science'
    normalized_text = re.sub(r"\bor\s+g/", "org/", normalized_text)
    # Step 4: remove spaces immediately after '?' so '? via=ihub' -> '?via=ihub'
    normalized_text = re.sub(r"\?\s+", "?", normalized_text)

    for m in url_pattern.finditer(normalized_text):
        url = m.group(0)
        cleaned_match = url.strip()
        url_status_ok = _url_is_ok(cleaned_match)
        key = ("URL", cleaned_match)
        if key in seen_pairs:
            continue
        seen_pairs.add(key)
        matches.append(
            {
                "label": "URL",
                "match": cleaned_match,
                "start": m.start(),
                "end": m.end(),
                "status_ok": url_status_ok,
            }
        )

    total = len(matches)
    if total == 0:
        msg_html = (
            '<div style="color: #666; padding: 10px; background-color: #f5f5f5; '
            'border-radius: 8px;">'
            "No explicit citation patterns (e.g., APA, numeric [1], DOI/URL) were clearly detected. "
            "Consider double-checking that your references are properly formatted."
            "</div>"
        )
        summary = "**Citations Detected:** 0"
        return msg_html, summary

    # Aggregate counts per pattern type
    from collections import Counter

    type_counts = Counter(m["label"] for m in matches)

    # Build HTML report
    rows = []
    # Limit to first 30 matches to keep UI manageable
    for idx, m in enumerate(matches[:30], start=1):
        match_text = html.escape(m["match"])
        # For URL citations, show status icon instead of context
        if m["label"] == "URL":
            if m.get("status_ok"):
                status_html = '<span style="color: #2e7d32; font-weight: bold;">✔</span>'
            else:
                status_html = '<span style="color: #c62828; font-weight: bold;">✖</span>'
        else:
            status_html = '<span style="color: #666;">–</span>'
        rows.append(
            f"""
            <tr>
                <td style="padding: 6px; border: 1px solid #ddd;">{idx}</td>
                <td style="padding: 6px; border: 1px solid #ddd;">{m['label']}</td>
                <td style="padding: 6px; border: 1px solid #ddd; font-family: monospace;">{match_text}</td>
                <td style="padding: 6px; border: 1px solid #ddd; text-align: center;">{status_html}</td>
            </tr>
            """
        )

    type_summary_html = "".join(
        f"<li><strong>{html.escape(label)}:</strong> {count}</li>"
        for label, count in type_counts.items()
    )

    table_html = "".join(rows)

    msg_html = f"""
    <div style="padding: 15px; background-color: #f5f5f5; border-radius: 8px;">
        <div style="margin-bottom: 10px;">
            <strong style="font-size: 16px;">Citation Patterns Detected</strong><br/>
            <span style="font-size: 13px; color: #555;">
                Found {total} potential citation instance(s) in your text.
            </span>
        </div>
        <ul style="font-size: 13px; color: #444; margin-bottom: 15px;">
            {type_summary_html}
        </ul>
        <div style="max-height: 420px; overflow-y: auto; background: white; border-radius: 6px; border: 1px solid #ddd;">
            <table style="width: 100%; border-collapse: collapse; font-size: 12px;">
                <thead style="background-color: #1976D2; color: white;">
                    <tr>
                        <th style="padding: 6px; border: 1px solid #ddd; text-align: left;">#</th>
                        <th style="padding: 6px; border: 1px solid #ddd; text-align: left;">Type</th>
                        <th style="padding: 6px; border: 1px solid #ddd; text-align: left;">Match</th>
                        <th style="padding: 6px; border: 1px solid #ddd; text-align: center;">URL Status</th>
                    </tr>
                </thead>
                <tbody>
                    {table_html}
                </tbody>
            </table>
        </div>
        <div style="margin-top: 10px; font-size: 12px; color: #777;">
            Note: This tool uses pattern matching and may miss some citations or flag non-citations that resemble references.
        </div>
    </div>
    """

    type_summary_md = ", ".join(f"{label}: {count}" for label, count in type_counts.items())
    summary = f"**Citations Detected:** {total} ({type_summary_md})"
    return msg_html, summary


def classify_text(text, progress=None, excluded_indices=None, return_debug_info=False, include_grammar_check=False):
    """
    Classifies the text and returns the result message and highlighted text.
    Plot generation code is commented out but can be re-enabled if needed.
    Args:
        text: Input text to classify
        progress: Progress callback
        excluded_indices: List of model indices to exclude (default: [6, 8, 9, 10] - human is included by default)
        return_debug_info: If True, also returns debug information
        include_grammar_check: If True, adds grammar error highlights (default: False - only for comprehensive scans)
    Returns: (result_message, highlighted_text) or (result_message, highlighted_text, debug_info)
    """
    # Default excluded indices (bloomz, davinci, dolly, dolly-v2-12b) - human is included by default
    if excluded_indices is None:
        excluded_indices = [6, 8, 9, 10]
    
    cleaned_text = clean_text(text)
    # If input is empty, clear the outputs
    if not cleaned_text.strip():
        if return_debug_info:
            return "", "", ""
        return "", ""

    if progress:
        progress(0.3, desc="Running AI models...")
    # Tokenize input and move to the appropriate device
    inputs = tokenizer(cleaned_text, return_tensors="pt", truncation=True, padding=True).to(device)

    # Perform inference with the three models
    with torch.no_grad():
        logits_1 = model_1(**inputs).logits
        logits_2 = model_2(**inputs).logits
        logits_3 = model_3(**inputs).logits

        softmax_1 = torch.softmax(logits_1, dim=1)
        softmax_2 = torch.softmax(logits_2, dim=1)
        softmax_3 = torch.softmax(logits_3, dim=1)

        averaged_probabilities = (softmax_1 + softmax_2 + softmax_3) / 3
        probabilities = averaged_probabilities[0]

    # Find the highest probability (excluding specified models)
    probs_for_decision = probabilities.clone()
    for idx in excluded_indices:
        if 0 <= idx < len(probs_for_decision):
            probs_for_decision[idx] = 0
    max_index = torch.argmax(probs_for_decision).item()
    max_prob = probabilities[max_index].item()
    
    # Get human probability (index 24) - set to 0 if human is excluded
    if 24 in excluded_indices:
        human_prob = 0.0
    else:
        human_prob = probabilities[24].item()
    
    # Get max AI probability (excluding human and excluded indices)
    ai_probs_for_max = probabilities.clone()
    for idx in excluded_indices:
        if 0 <= idx < len(ai_probs_for_max):
            ai_probs_for_max[idx] = 0
    # Also exclude human (24) from AI max calculation
    if 24 not in excluded_indices:
        ai_probs_for_max[24] = 0
    # Don't filter by 0.02 threshold - use the actual max probability from included models
    max_ai_index = torch.argmax(ai_probs_for_max).item()
    max_ai_prob = probabilities[max_ai_index].item()
    
    # Calculate initial percentages for highlighting (based on document-level probabilities)
    # Compare human vs max AI to determine the correct highlighting percentage
    total_prob_for_detection = human_prob + max_ai_prob
    if total_prob_for_detection > 0:
        # Calculate AI percentage - if human is higher, AI percentage will be lower (correct)
        ai_percentage_for_highlighting = (max_ai_prob / total_prob_for_detection) * 100
        # Ensure it's between 0 and 100
        ai_percentage_for_highlighting = max(0.0, min(100.0, ai_percentage_for_highlighting))
    else:
        # If no probabilities, default to 50/50
        ai_percentage_for_highlighting = 50.0
    
    # If human_prob is higher than max_ai_prob, ensure AI percentage reflects that
    # This is critical - when human is higher, most text should be highlighted as human
    if human_prob > max_ai_prob:
        # Recalculate AI percentage - it should be less than 50% when human is higher
        calculated_ai_pct = (max_ai_prob / total_prob_for_detection) * 100
        # Force it to be less than 50% to ensure more text is highlighted as human
        if calculated_ai_pct >= 50.0:
            # This shouldn't happen mathematically, but if it does, cap at 49%
            ai_percentage_for_highlighting = 49.0
        else:
            ai_percentage_for_highlighting = calculated_ai_pct
        # Double-check: if human is higher, AI percentage must be < 50%
        assert ai_percentage_for_highlighting < 50.0, f"Error: human_prob ({human_prob}) > max_ai_prob ({max_ai_prob}) but ai_percentage is {ai_percentage_for_highlighting}"
    
    # Generate highlighted text first - this will determine the actual detection
    if return_debug_info:
        highlighted_text, ai_char_count, human_char_count, sentence_debug_data = highlight_ai_text(
            text, ai_percentage_for_highlighting, progress, excluded_indices, return_sentence_data=True
        )
        sentence_data, sentence_ai_percentages, ai_indices, overall_ai_pct, sentence_texts = sentence_debug_data
    else:
        highlighted_text, ai_char_count, human_char_count = highlight_ai_text(
            text, ai_percentage_for_highlighting, progress, excluded_indices, return_sentence_data=False
        )
        sentence_data = []
        sentence_ai_percentages = []
        ai_indices = set()
        overall_ai_pct = ai_percentage_for_highlighting
        sentence_texts = []
    
    # Add grammar highlights to the highlighted text (only if requested - for comprehensive scans)
    if include_grammar_check:
        if progress:
            progress(0.85, desc="Checking grammar...")
        highlighted_text, grammar_error_count = add_grammar_highlights_to_text(text, highlighted_text)
    else:
        grammar_error_count = 0
    
    if return_debug_info:
        debug_info = ""
    else:
        debug_info = ""
    
    # Calculate percentages based on actual highlighted text area (this is what determines detection)
    total_chars = ai_char_count + human_char_count
    if total_chars > 0:
        human_area_percentage_raw = (human_char_count / total_chars) * 100
        ai_area_percentage_raw = (ai_char_count / total_chars) * 100
    else:
        human_area_percentage_raw = 0.0
        ai_area_percentage_raw = 0.0
    
    # Round human percentage up to next integer, AI is the remainder
    human_detection_percentage = math.ceil(human_area_percentage_raw)
    ai_detection_percentage = 100 - human_detection_percentage
    
    # Determine detected model based on highlighted text ratio
    # This ensures detection matches what's actually highlighted
    if human_area_percentage_raw > ai_area_percentage_raw:
        detected_model = "Human"
        model_color = "#4CAF50"  # Green for Human
    elif ai_area_percentage_raw > human_area_percentage_raw:
        detected_model = "AI"
        model_color = "#FF5733"  # Red for AI
    else:
        # If equal, check document-level probabilities as tiebreaker
        if human_prob >= max_ai_prob:
            detected_model = "Human"
            model_color = "#4CAF50"  # Green for Human
        else:
            detected_model = "AI"
            model_color = "#FF5733"  # Red for AI
    
    # Use highlighted text area percentages for the progress bar and detection
    # This ensures everything is based on the actual highlighted text ratio
    result_message = f"""
    <div style="width: 100%; margin: 20px 0;">
        <div style="display: flex; justify-content: space-between; margin-bottom: 8px; font-size: 14px;">
            <span style="color: #4CAF50; font-weight: bold;">Human: {human_detection_percentage:.1f}%</span>
            <span style="color: #FF5733; font-weight: bold;">AI: {ai_detection_percentage:.1f}%</span>
        </div>
        <div style="width: 100%; height: 30px; background-color: #e0e0e0; border-radius: 15px; overflow: hidden; display: flex;">
            <div style="width: {human_detection_percentage}%; background-color: #4CAF50; transition: width 0.3s ease;"></div>
            <div style="width: {ai_detection_percentage}%; background-color: #FF5733; transition: width 0.3s ease;"></div>
        </div>
        <div style="margin-top: 12px; text-align: center; font-size: 14px;">
            <span style="color: #666; font-weight: bold;">Detected: </span>
            <span style="color: {model_color}; font-weight: bold; font-size: 16px;">{detected_model}</span>
        </div>
    </div>
    """


    # Generate debug info if requested
    if return_debug_info:
        # Generate per-sentence debug info using current calculation
        if sentence_data:
            debug_info = generate_per_sentence_debug_info(
                sentence_data, sentence_ai_percentages, ai_indices, overall_ai_pct, excluded_indices, sentence_texts
            )
        else:
            # Fallback to document-level debug info if no sentence data
            # Get all probabilities with their model names
            prob_list = []
            for idx in range(len(probabilities)):
                if idx < len(label_mapping):
                    model_name = label_mapping[idx]
                    prob = probabilities[idx].item()
                    prob_list.append((idx, model_name, prob))
            
            # Sort by probability (descending)
            prob_list.sort(key=lambda x: x[2], reverse=True)
            
            # Calculate total probability for percentage calculation
            total_prob = sum(p[2] for p in prob_list)
            
            # Generate debug table HTML
            debug_info = """
            <div style="margin-top: 20px; padding: 15px; background-color: #f5f5f5; border-radius: 8px; border: 1px solid #ddd;">
                <h3 style="margin-top: 0; color: #333;">Debug Information - All Models</h3>
                <div style="overflow-x: auto;">
                    <table style="width: 100%; border-collapse: collapse; font-family: monospace; font-size: 12px;">
                        <thead>
                            <tr style="background-color: #4CAF50; color: white;">
                                <th style="padding: 8px; text-align: left; border: 1px solid #ddd;">Rank</th>
                                <th style="padding: 8px; text-align: left; border: 1px solid #ddd;">Model</th>
                                <th style="padding: 8px; text-align: right; border: 1px solid #ddd;">Probability</th>
                                <th style="padding: 8px; text-align: right; border: 1px solid #ddd;">Percentage</th>
                                <th style="padding: 8px; text-align: center; border: 1px solid #ddd;">Excluded</th>
                            </tr>
                        </thead>
                        <tbody>
            """
            
            for rank, (idx, model_name, prob) in enumerate(prob_list, 1):
                percentage = (prob / total_prob * 100) if total_prob > 0 else 0.0
                # Always round upward to next whole number if decimal exists
                if percentage > 0 and percentage % 1 != 0:
                    percentage = math.ceil(percentage)
                else:
                    percentage = int(percentage)
                is_excluded = idx in excluded_indices
                excluded_marker = "✓" if is_excluded else ""
                row_color = "#ffebee" if is_excluded else "#ffffff"
                
                debug_info += f"""
                            <tr style="background-color: {row_color};">
                                <td style="padding: 6px; border: 1px solid #ddd;">{rank}</td>
                                <td style="padding: 6px; border: 1px solid #ddd; font-weight: {'bold' if idx == max_index or idx == max_ai_index else 'normal'}">{model_name}</td>
                                <td style="padding: 6px; border: 1px solid #ddd; text-align: right;">{prob:.6f}</td>
                                <td style="padding: 6px; border: 1px solid #ddd; text-align: right;">{percentage}%</td>
                                <td style="padding: 6px; border: 1px solid #ddd; text-align: center; color: {'#d32f2f' if is_excluded else '#666'}">{excluded_marker}</td>
                            </tr>
                """
            
            debug_info += """
                        </tbody>
                    </table>
                </div>
                <p style="margin-top: 10px; font-size: 11px; color: #666;">
                    <strong>Note:</strong> Excluded models are marked with ✓ and highlighted in red. 
                    Models with bold names are the detected max probability models.
                </p>
            </div>
            """
    
    # Plot generation code (commented out - can be re-enabled if needed)
    # ai_probs_for_plot = probabilities.clone()
    # ai_probs_for_plot[6] = 0   # Exclude bloomz
    # ai_probs_for_plot[8] = 0   # Exclude davinci
    # ai_probs_for_plot[9] = 0   # Exclude dolly
    # ai_probs_for_plot[10] = 0  # Exclude dolly-v2-12b
    # top_5_probs, top_5_indices = torch.topk(ai_probs_for_plot, 5)
    #
    # top_5_probs = top_5_probs.cpu().numpy()
    # top_5_labels = [label_mapping[i.item()] if i.item() < len(label_mapping) and label_mapping[i.item()] is not None else f'Unknown (index {i.item()})' for i in top_5_indices]
    #
    # fig, ax = plt.subplots(figsize=(10, 5))
    # bars = ax.barh(top_5_labels, top_5_probs, color='#4CAF50', alpha=0.8)
    # ax.set_xlabel('Probability', fontsize=12)
    # ax.set_title('Top 5 Predictions', fontsize=14, fontweight='bold')
    # ax.invert_yaxis()
    # ax.grid(axis='x', linestyle='--', alpha=0.6)
    #
    # for bar in bars:
    #     width = bar.get_width()
    #     label_x_pos = width + 0.01
    #     ax.text(label_x_pos, bar.get_y() + bar.get_height() / 2, f'{width:.2%}', va='center')
    #
    # ax.set_xlim(0, max(top_5_probs) * 1.18) 
    # plt.tight_layout()
    # fig = fig  # Keep fig for return statement
    
    
    if return_debug_info:
        cleanup_memory()
        return result_message, highlighted_text, debug_info
    
    cleanup_memory()
    return result_message, highlighted_text  # , fig  # Uncomment fig to re-enable plot

def extract_and_populate_text(file):
    """Extract text from uploaded file and return it to populate the textarea"""
    if file is None or file == "":
        return ""
    
    file_path = file if isinstance(file, str) else (file.name if hasattr(file, 'name') else str(file))
    if not file_path or file_path == "":
        return ""
    
    # Save uploaded file to uploads directory
    saved_path = save_uploaded_file(file_path)
    
    extracted_text = extract_text_from_file(saved_path)
    
    # Don't put error messages in textarea - return empty string for errors
    if extracted_text.startswith("Error") or extracted_text.startswith("Unsupported"):
        return ""
    
    return extracted_text if extracted_text else ""

def check_file_error(file):
    """Check if file extraction resulted in an error and return error message"""
    if file is None or file == "":
        return ""
    
    file_path = file if isinstance(file, str) else (file.name if hasattr(file, 'name') else str(file))
    if not file_path or file_path == "":
        return ""
    
    # Save uploaded file to uploads directory
    saved_path = save_uploaded_file(file_path)
    
    extracted_text = extract_text_from_file(saved_path)
    
    if extracted_text.startswith("Error") or extracted_text.startswith("Unsupported"):
        error_msg = f'<div style="color: red; font-weight: bold;">{html.escape(extracted_text)}</div>'
        return error_msg
    
    return ""

def process_file_or_text(file, text_input, progress=gr.Progress(), excluded_indices=None, return_debug_info=False, include_grammar_check=False):
    """Process either uploaded file or text input"""
    # Show initial progress to make progress bar visible
    if progress:
        progress(0.0, desc="Starting analysis...")
    
    # Prioritize text_input if it exists (it will be populated when file is uploaded)
    if text_input and text_input.strip():
        # Check language first - only English is supported
        is_valid, error_msg = check_language_and_return_error(text_input)
        if not is_valid:
            if return_debug_info:
                return text_input, error_msg, "", ""
            return text_input, error_msg, ""
        
        # Use text input for classification
        if return_debug_info:
            result_msg, highlighted, debug_info = classify_text(text_input, progress, excluded_indices, return_debug_info=True, include_grammar_check=include_grammar_check)
            if progress:
                progress(1.0, desc="Analysis complete!")
            return text_input, result_msg, highlighted, debug_info
        else:
            result_msg, highlighted = classify_text(text_input, progress, excluded_indices, return_debug_info=False, include_grammar_check=include_grammar_check)
            if progress:
                progress(1.0, desc="Analysis complete!")
            # Return highlighted text to be shown in the text area, and result message
            return text_input, result_msg, highlighted
    elif file is not None and file != "":
        # Extract text from uploaded file (fallback if text_input is empty)
        file_path = file if isinstance(file, str) else (file.name if hasattr(file, 'name') else str(file))
        if not file_path or file_path == "":
            if return_debug_info:
                return "", '<div style="color: red; font-weight: bold;">Please upload a file or enter text.</div>', "", ""
            return "", '<div style="color: red; font-weight: bold;">Please upload a file or enter text.</div>', ""
        
        # Save uploaded file to uploads directory
        saved_path = save_uploaded_file(file_path)
        
        progress(0.1, desc="Extracting text from file...")
        print(f"Processing file: {saved_path}")  # Debug output
        extracted_text = extract_text_from_file(saved_path)
        
        if extracted_text.startswith("Error") or extracted_text.startswith("Unsupported"):
            escaped_error = html.escape(extracted_text)
            if return_debug_info:
                return "", f'<div style="color: red; font-weight: bold;">{escaped_error}</div>', "", ""
            return "", f'<div style="color: red; font-weight: bold;">{escaped_error}</div>', ""
        
        if not extracted_text or not extracted_text.strip():
            if return_debug_info:
                return "", '<div style="color: red; font-weight: bold;">Error: No text could be extracted from the file.</div>', "", ""
            return "", '<div style="color: red; font-weight: bold;">Error: No text could be extracted from the file.</div>', ""
        
        print(f"Extracted {len(extracted_text)} characters from file")  # Debug output
        
        # Check language first - only English is supported
        is_valid, error_msg = check_language_and_return_error(extracted_text)
        if not is_valid:
            if return_debug_info:
                return extracted_text, error_msg, "", ""
            return extracted_text, error_msg, ""
        
        # Use extracted text for classification
        if return_debug_info:
            result_msg, highlighted, debug_info = classify_text(extracted_text, progress, excluded_indices, return_debug_info=True, include_grammar_check=include_grammar_check)
            if progress:
                progress(1.0, desc="Analysis complete!")
            return extracted_text, result_msg, highlighted, debug_info
        else:
            result_msg, highlighted = classify_text(extracted_text, progress, excluded_indices, return_debug_info=False, include_grammar_check=include_grammar_check)
            if progress:
                progress(1.0, desc="Analysis complete!")
            # Return extracted text, result message, and highlighted version
            return extracted_text, result_msg, highlighted
    else:
        if return_debug_info:
            return "", "", "", ""
        return "", "", ""



title = "AI Content Detector & Writing Assistant"

description = """
Comprehensive writing analysis tool: AI Detection, Grammar Check, TOEFL Rubric, Writing Feedback, and Plagiarism Detection.

**Supported file formats:** PDF, DOCX, PPTX, XLSX, CSV, TXT, MD, RTF
"""
bottom_text = "**Developed by Arekiv**"

AI_texts = [
"Camels are remarkable desert animals known for their unique adaptations to harsh, arid environments. Native to the Middle East, North Africa, and parts of Asia, camels have been essential to human life for centuries, serving as a mode of transportation, a source of food, and even a symbol of endurance and survival. There are two primary species of camels: the dromedary camel, which has a single hump and is commonly found in the Middle East and North Africa, and the Bactrian camel, which has two humps and is native to Central Asia. Their humps store fat, not water, as commonly believed, allowing them to survive long periods without food by metabolizing the stored fat for energy. Camels are highly adapted to desert life. They can go for weeks without water, and when they do drink, they can consume up to 40 gallons in one sitting. Their thick eyelashes, sealable nostrils, and wide, padded feet protect them from sand and help them walk easily on loose desert terrain.",
"Wines are a fascinating reflection of culture, history, and craftsmanship. They embody a rich diversity shaped by the land, climate, and traditions where they are produced. From the bold reds of Bordeaux to the crisp whites of New Zealand, each bottle tells a unique story. What makes wine so special is its ability to connect people. Whether shared at a family dinner, a celebratory event, or a quiet evening with friends, wine enhances experiences and brings people together. The variety of flavors and aromas, influenced by grape type, fermentation techniques, and aging processes, make wine tasting a complex yet rewarding journey for the senses.",
"I find artificial intelligence (AI) to be one of the most transformative and fascinating technologies of our time. Its potential spans a wide range of applications, from automating mundane tasks to revolutionizing industries like healthcare, education, and entertainment. AI has already made significant contributions in fields like language processing, image recognition, and decision-making systems, enabling innovations that were once purely science fiction. However, as powerful as AI can be, it also brings challenges and responsibilities. Ethical considerations, such as bias in data, transparency, and the potential for misuse, need to be carefully addressed to ensure fairness and accountability. The rise of generative AI has also sparked debates about creativity, originality, and intellectual property, making it essential to strike a balance between technological advancement and respecting human contributions."
]

Human_texts = [
"The present book is intended as a text in basic mathematics. As such, it can have multiple use: for a one-year course in the high schools during the third or fourth year (if possible the third, so that calculus can be taken during the fourth year); for a complementary reference in earlier high school grades (elementary algebra and geometry are covered); for a one-semester course at the college level, to review or to get a firm foundation in the basic mathematics necessary to go ahead in calculus, linear algebra, or other topics. Years ago, the colleges used to give courses in “ college algebra” and other subjects which should have been covered in high school. More recently, such courses have been thought unnecessary, but some experiences I have had show that they are just as necessary as ever. What is happening is that thecolleges are getting a wide variety of students from high schools, ranging from exceedingly well-prepared ones who have had a good first course in calculus, down to very poorly prepared ones.",
"Fats are rich in energy, build body cells, support brain development of infants, help body processes, and facilitate the absorption and use of fat-soluble vitamins A, D, E, and K. The major component of lipids is glycerol and fatty acids. According to chemical properties, fatty acids can be divided into saturated and unsaturated fatty acids. Generally lipids containing saturated fatty acids are solid at room temperature and include animal fats (butter, lard, tallow, ghee) and tropical oils (palm,coconut, palm kernel). Saturated fats increase the risk of heart disease.",
"To make BERT handle a variety of down-stream tasks, our input representation is able to unambiguously represent both a single sentence and a pair of sentences (e.g., h Question, Answeri) in one token sequence. Throughout this work, a “sentence” can be an arbitrary span of contiguous text, rather than an actual linguistic sentence. A “sequence” refers to the input token sequence to BERT, which may be a single sentence or two sentences packed together. We use WordPiece embeddings (Wu et al., 2016) with a 30,000 token vocabulary. The first token of every sequence is always a special classification token ([CLS]). The final hidden state corresponding to this token is used as the aggregate sequence representation for classification tasks. Sentence pairs are packed together into a single sequence."]

# Custom CSS styles
custom_css = """
    .scrollable-debug {
        max-height: 1750px !important;
        overflow-y: auto !important;
        overflow-x: hidden !important;
    }
    .scrollable-debug > div {
        max-height: 1750px !important;
        overflow-y: auto !important;
        overflow-x: hidden !important;
    }
    @import url('https://fonts.googleapis.com/css2?family=Roboto+Mono:wght@400;700&display=swap');
    #text_input_box { border-radius: 10px; border: 2px solid #4CAF50; font-size: 18px; padding: 15px; margin-bottom: 20px; width: 100%; box-sizing: border-box; margin: auto; }
    #highlighted_text_display { border-radius: 10px; border: 2px solid #4CAF50; font-size: 18px; padding: 15px; margin-bottom: 20px; width: 100%; box-sizing: border-box; margin: auto; min-height: 75px; line-height: 1.6; cursor: pointer; background-color: white; font-family: 'Roboto Mono', monospace; white-space: pre-wrap; }
    #highlighted_text_display:hover { border-color: #45a049; box-shadow: 0 0 5px rgba(76, 175, 80, 0.3); }
    #highlighted_text_display > div { max-height: 350px; overflow-y: auto !important; overflow-x: hidden !important; padding: 0; }
    @media (max-width: 768px) { #highlighted_text_display { width: 100%; } }
    .form.svelte-633qhp { background: none; border: none; box-shadow: none; }
    #result_output_box { border-radius: 10px; border: 2px solid #4CAF50; font-size: 18px; padding: 15px; margin-top: 20px; width: 90%; box-sizing: border-box; text-align: center; margin: auto; }
    @media (max-width: 768px) { #result_output_box { width: 100%; } #text_input_box{ width: 100%; } }
    body { font-family: 'Roboto Mono', sans-serif !important; padding: 20px; min-height: 100vh; }
    html { overflow-y: scroll; }
    .gradio-container { border: 1px solid #4CAF50; border-radius: 15px; padding: 30px; box-shadow: 0px 0px 10px rgba(0,255,0,0.6); width: 100%; margin: 0; }
    h1 { text-align: center; font-size: 32px; font-weight: bold; margin-bottom: 30px; }
    .highlight-human { color: #4CAF50; font-weight: bold; background: rgba(76, 175, 80, 0.2); padding: 5px; border-radius: 8px; }
    .highlight-ai { color: #FF5733; font-weight: bold; background: rgba(255, 87, 51, 0.2); padding: 5px; border-radius: 8px; }
    .ai-highlight { background-color: #ffebee; color: #c62828; padding: 3px 5px; border-radius: 4px; border-left: 3px solid #d32f2f; font-weight: 500; box-shadow: 0 1px 2px rgba(211, 47, 47, 0.1); }
    .human-highlight { background-color: #e8f5e9; color: #000000; padding: 3px 5px; border-radius: 4px; font-weight: 500; box-shadow: 0 1px 2px rgba(56, 142, 60, 0.1); }
    #highlighted_text_box { border-radius: 10px; border: 2px solid #4CAF50; font-size: 16px; padding: 15px; margin-top: 20px; width: 90%; box-sizing: border-box; margin: auto; line-height: 1.6; overflow: visible; }
    #highlighted_text_box > div { max-height: 600px; overflow-y: auto !important; overflow-x: hidden !important; }
    #bottom_text { text-align: center; margin-top: 50px; font-weight: bold; font-size: 20px; }
    .grammar-summary { margin-bottom: 10px; display: block; font-weight: bold; color: #333; }
    .grammar-text { white-space: pre-wrap; font-family: inherit; }
    .grammar-error-word { color: inherit; text-decoration: underline wavy #ef4444; text-decoration-thickness: 2px; text-underline-offset: 2px; cursor: pointer; display: inline; white-space: nowrap; }
    .ai-highlight .grammar-error-word { text-decoration: underline wavy #dc2626; color: inherit; }
    .human-highlight .grammar-error-word { text-decoration: underline wavy #ef4444; color: inherit; }
    .grammar-error-underline { background-color: #fee2e2; color: #991b1b; padding: 2px 4px; border-radius: 3px; border: 1px solid #fca5a5; cursor: pointer; display: inline; white-space: nowrap; font-weight: 500; }
    .grammar-suggestion-popup { 
        position: absolute; 
        background-color: #fff; 
        border: 2px solid #FF5733; 
        border-radius: 8px; 
        padding: 12px; 
        box-shadow: 0 4px 12px rgba(0,0,0,0.15); 
        z-index: 1000; 
        max-width: 400px; 
        font-size: 14px; 
        line-height: 1.5;
        display: none;
    }
    .grammar-suggestion-popup.show { display: block; }
    .grammar-suggestion-popup .message { 
        font-weight: bold; 
        color: #FF5733; 
        margin-bottom: 8px; 
        border-bottom: 1px solid #eee; 
        padding-bottom: 6px;
    }
    .grammar-suggestion-popup .suggestions { 
        margin-top: 8px; 
    }
    .grammar-suggestion-popup .suggestion-item { 
        padding: 6px 10px; 
        margin: 4px 0; 
        background-color: #f5f5f5; 
        border-radius: 4px; 
        cursor: pointer; 
        transition: background-color 0.2s;
    }
    .grammar-suggestion-popup .suggestion-item:hover { 
        background-color: #e0e0e0; 
    }
    .grammar-suggestion-popup .close-btn { 
        float: right; 
        cursor: pointer; 
        font-weight: bold; 
        color: #666; 
        font-size: 18px; 
        line-height: 1;
    }
    .grammar-suggestion-popup .close-btn:hover { 
        color: #FF5733; 
    }
    .block.svelte-11xb1hd{ background: none !important; }
"""

iface = gr.Blocks()

with iface:
    # Inject custom CSS and JavaScript via HTML component (compatible with all Gradio versions)
    gr.HTML(f"""
    <style>
    {custom_css}
    </style>
    <script>
    (function() {{
        function showGrammarSuggestion(element) {{
            console.log('Grammar error clicked!', element);
            // Remove any existing popups
            const existingPopup = document.querySelector('.grammar-suggestion-popup');
            if (existingPopup) {{
                existingPopup.remove();
            }}
            
            // Get error data from element
            const message = element.getAttribute('data-grammar-message') || 'Grammar error';
            const suggestionsStr = element.getAttribute('data-grammar-suggestions') || '';
            const suggestions = suggestionsStr ? suggestionsStr.split('|') : [];
            
            console.log('Message:', message, 'Suggestions:', suggestions);
            
            // Create popup element
            const popup = document.createElement('div');
            popup.className = 'grammar-suggestion-popup show';
            
            // Get position relative to viewport
            const rect = element.getBoundingClientRect();
            const scrollTop = window.pageYOffset || document.documentElement.scrollTop;
            const scrollLeft = window.pageXOffset || document.documentElement.scrollLeft;
            
            // Position popup above or below the element
            let top = rect.top + scrollTop - 10;
            let left = rect.left + scrollLeft;
            
            // Adjust if popup would go off screen
            if (top < scrollTop) {{
                top = rect.bottom + scrollTop + 10;
            }}
            if (left + 400 > window.innerWidth + scrollLeft) {{
                left = window.innerWidth + scrollLeft - 420;
            }}
            if (left < scrollLeft) {{
                left = scrollLeft + 10;
            }}
            
            popup.style.top = top + 'px';
            popup.style.left = left + 'px';
            popup.style.position = 'absolute';
            popup.style.zIndex = '10000';
            
            // Build popup content
            let content = '<span class="close-btn">&times;</span>';
            content += '<div class="message">' + message + '</div>';
            
            if (suggestions.length > 0) {{
                content += '<div class="suggestions"><strong>Suggestions:</strong>';
                suggestions.forEach(function(suggestion, index) {{
                    content += '<div class="suggestion-item" data-index="' + index + '">' + 
                              suggestion + '</div>';
                }});
                content += '</div>';
            }} else {{
                content += '<div class="suggestions">No suggestions available.</div>';
            }}
            
            // Store suggestions in popup for replacement function
            popup.setAttribute('data-suggestions', JSON.stringify(suggestions));
            
            popup.innerHTML = content;
            document.body.appendChild(popup);
            console.log('Popup created and added to DOM');
        }}
        
        function replaceGrammarError(element, index) {{
            const popup = element.closest('.grammar-suggestion-popup');
            if (!popup) return;
            
            const suggestions = JSON.parse(popup.getAttribute('data-suggestions') || '[]');
            if (index >= 0 && index < suggestions.length) {{
                const replacement = suggestions[index];
                // Copy to clipboard
                if (navigator.clipboard && navigator.clipboard.writeText) {{
                    navigator.clipboard.writeText(replacement).then(function() {{
                        // Show temporary feedback
                        element.style.backgroundColor = '#4CAF50';
                        element.style.color = 'white';
                        setTimeout(function() {{
                            element.style.backgroundColor = '';
                            element.style.color = '';
                        }}, 1000);
                    }}).catch(function() {{
                        // Fallback: show alert
                        alert('Suggestion: ' + replacement + '\\n\\nPlease copy this and replace the error manually.');
                    }});
                }} else {{
                    // Fallback for browsers without clipboard API
                    alert('Suggestion: ' + replacement + '\\n\\nPlease copy this and replace the error manually.');
                }}
            }}
            // Close popup after a short delay
            setTimeout(function() {{
                if (popup) popup.remove();
            }}, 1500);
        }}
        
        // Setup event delegation - use capture phase to catch events early
        function setupEventListeners() {{
            console.log('Setting up event listeners');
            // Remove old listeners if any
            document.removeEventListener('click', handleClick, true);
            document.addEventListener('click', handleClick, true);
        }}
        
        function handleClick(e) {{
            // Check if clicked element is a grammar error or inside one
            const grammarError = e.target.closest('.grammar-error-underline');
            if (grammarError) {{
                console.log('Grammar error span clicked!');
                e.stopPropagation();
                e.preventDefault();
                showGrammarSuggestion(grammarError);
                return false;
            }}
            
            // Check if clicked on suggestion item
            const suggestionItem = e.target.closest('.suggestion-item');
            if (suggestionItem) {{
                e.stopPropagation();
                e.preventDefault();
                const index = parseInt(suggestionItem.getAttribute('data-index') || '0');
                replaceGrammarError(suggestionItem, index);
                return false;
            }}
            
            // Check if clicked on close button
            if (e.target.classList.contains('close-btn')) {{
                e.stopPropagation();
                e.preventDefault();
                const popup = e.target.closest('.grammar-suggestion-popup');
                if (popup) popup.remove();
                return false;
            }}
            
            // Close popup if clicking outside
            const popup = document.querySelector('.grammar-suggestion-popup');
            if (popup && !popup.contains(e.target)) {{
                popup.remove();
            }}
        }}
        
        // Initialize immediately
        setupEventListeners();
        
        // Also initialize on page load
        if (document.readyState === 'loading') {{
            document.addEventListener('DOMContentLoaded', setupEventListeners);
        }} else {{
            setupEventListeners();
        }}
        
        // Re-setup after Gradio updates content (using MutationObserver)
        const observer = new MutationObserver(function(mutations) {{
            let shouldReattach = false;
            mutations.forEach(function(mutation) {{
                if (mutation.addedNodes.length > 0) {{
                    mutation.addedNodes.forEach(function(node) {{
                        if (node.nodeType === 1) {{ // Element node
                            if (node.classList && node.classList.contains('grammar-error-underline')) {{
                                shouldReattach = true;
                            }}
                            if (node.querySelector && node.querySelector('.grammar-error-underline')) {{
                                shouldReattach = true;
                            }}
                        }}
                    }});
                }}
            }});
            if (shouldReattach) {{
                console.log('New grammar errors detected, reattaching listeners');
                setTimeout(setupEventListeners, 100);
            }}
        }});
        
        // Observe changes to the document body
        if (document.body) {{
            observer.observe(document.body, {{
                childList: true,
                subtree: true
            }});
        }} else {{
            document.addEventListener('DOMContentLoaded', function() {{
                observer.observe(document.body, {{
                    childList: true,
                    subtree: true
                }});
            }});
        }}
        
        // Also setup when window loads
        window.addEventListener('load', setupEventListeners);
        
        // Make function globally accessible
        window.showGrammarSuggestion = showGrammarSuggestion;
    }})();
    </script>
    """, visible=False)
    gr.Markdown(f"# {title}")
    gr.Markdown(description)
    
    # Shared file upload and text input (used by all tabs)
    with gr.Row():
        file_upload = gr.File(
            label="Upload Document",
            file_types=[".pdf", ".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".csv", ".txt", ".md", ".rtf"],
            type="filepath"
        )
    
    # Shared text input for all tabs
    text_input = gr.Textbox(
        label="Enter text here (shared across all tabs)", 
        placeholder="Type or paste your content here...", 
        elem_id="text_input_box", 
        lines=10
    )
    
    # Common Scan button for all tabs (will change to Stop during processing)
    process_btn = gr.Button("Scan", variant="primary", size="lg")
    scanning_state = gr.State(value=False)  # Track if scanning is in progress
    
    # Create tabs for different features
    with gr.Tabs():
        # Tab 1: AI Detection (existing functionality)
        with gr.Tab("AI Detection"):
            # Grammar language selection for grammar checking (shown as red underlines in highlighted text)
            grammar_language = gr.Dropdown(
                label="Grammar Check Language",
                choices=[
                    ("English (US)", "en-US"),
                    ("English (UK)", "en-GB"),
                    ("English (CA)", "en-CA"),
                    ("English (AU)", "en-AU"),
                    ("German", "de-DE"),
                    ("French", "fr-FR"),
                    ("Spanish", "es-ES"),
                    ("Italian", "it-IT"),
                    ("Portuguese", "pt-PT"),
                    ("Portuguese (BR)", "pt-BR"),
                    ("Dutch", "nl-NL"),
                    ("Polish", "pl-PL"),
                    ("Russian", "ru-RU"),
                    ("Chinese", "zh"),
                    ("Japanese", "ja"),
                    ("Korean", "ko"),
                    ("Arabic", "ar"),
                    ("Hindi", "hi"),
                ],
                value="en-US",
                info="Grammar errors will be shown as red underlines in the highlighted text"
            )
            # State to store last processed text, result message, and highlighted result
            last_processed_text = gr.State(value="")
            last_result_msg = gr.State(value="")
            last_highlighted_html = gr.State(value="")
            
            result_output = gr.Markdown("", elem_id="result_output_box")
            highlighted_text_display = gr.HTML("", elem_id="highlighted_text_display", visible=True)
        
        # Tab 2: TOEFL Rubric
        with gr.Tab("TOEFL Rubric"):
            toefl_output = gr.HTML("")
            
            def process_toefl(text, progress=gr.Progress()):
                if progress:
                    progress(0.5, desc="Calculating TOEFL scores...")
                scores = calculate_toefl_score(text)
                if scores:
                    return format_toefl_rubric(scores, text)
                return ""
            
        
        # Tab 4: Writing Feedback
        with gr.Tab("Writing Feedback"):
            gr.Markdown("**Get comprehensive writing feedback** - Analyze your writing with detailed suggestions.")
            feedback_output = gr.HTML("")
            
            def process_feedback(text, progress=gr.Progress()):
                if progress:
                        progress(0.5, desc="Generating feedback...")
                result = generate_writing_feedback(text)
                if progress:
                    progress(1.0, desc="Feedback generated!")
                return result
        
        # Tab 5: Plagiarism Check
        with gr.Tab("Plagiarism Check"):
            gr.Markdown("**Powered by PrepostSEO Plagiarism Checker API** - Check your text against billions of web pages for plagiarism.")
            gr.Markdown("**Note:** Maximum word limit is 5,000 words. Text exceeding this limit will be automatically truncated.")
            exclude_urls_input = gr.Textbox(label="Exclude URLs (optional - one per line)", placeholder="Enter URLs to exclude from the plagiarism check (e.g., your own website)...", lines=3)
            plagiarism_output = gr.HTML("")
            plagiarism_score = gr.Markdown("")
            
            def process_plagiarism(text, exclude_urls, progress=gr.Progress()):
                if progress:
                    progress(0.3, desc="Sending request to plagiarism API...")
                exclude_list = [url.strip() for url in exclude_urls.split('\n') if url.strip()] if exclude_urls else None
                result_html, similarity = check_plagiarism(text, exclude_list)
                score_text = f"**Plagiarism Score:** {similarity:.1f}%"
                if progress:
                    progress(1.0, desc="Analysis complete!")
                return result_html, score_text
            
            exclude_urls_input.change(
                process_plagiarism, inputs=[text_input, exclude_urls_input], outputs=[plagiarism_output, plagiarism_score]
            )
        
        # Tab 6: Citation Checker
        with gr.Tab("Citation Checker"):
            gr.Markdown("**Check whether your writing includes clearly formatted citations** (APA-style, numeric [1], DOI/URL-based, etc.).")
            citation_output = gr.HTML("")
            citation_stats = gr.Markdown("")
        
        # Tab 7: Developer Options (only show if DEV_MODE is enabled)
        # Initialize dev components outside tab scope for Scan button access
        checkbox_inputs = []
        dev_result_output = None
        dev_debug_output = None
        
        if DEV_MODE:
            with gr.Tab("Developer Options"):
                gr.Markdown("**🔧 Developer Options** - Configure model exclusions and view detailed debug information.")
                
                with gr.Row():
                    with gr.Column(scale=1):
                        gr.Markdown("### Model Inclusion/Exclusion")
                        gr.Markdown("**Select which models to include in detection:** (Unchecked = Excluded)")
                        
                        # Default excluded models: bloomz (6), davinci (8), dolly (9), dolly-v2-12b (10) - human is included by default
                        default_excluded_models = [6, 8, 9, 10]
                        
                        # Create checkboxes in 3 columns
                        model_checkboxes = []
                        with gr.Row():
                            with gr.Column():
                                for idx in range(0, len(label_mapping), 3):
                                    if idx < len(label_mapping):
                                        is_checked = idx not in default_excluded_models
                                        checkbox = gr.Checkbox(
                                            label=f"[{idx}] {label_mapping[idx]}",
                                            value=is_checked,
                                            elem_id=f"model_checkbox_{idx}"
                                        )
                                        model_checkboxes.append((idx, checkbox))
                            
                            with gr.Column():
                                for idx in range(1, len(label_mapping), 3):
                                    if idx < len(label_mapping):
                                        is_checked = idx not in default_excluded_models
                                        checkbox = gr.Checkbox(
                                            label=f"[{idx}] {label_mapping[idx]}",
                                            value=is_checked,
                                            elem_id=f"model_checkbox_{idx}"
                                        )
                                        model_checkboxes.append((idx, checkbox))
                            
                            with gr.Column():
                                for idx in range(2, len(label_mapping), 3):
                                    if idx < len(label_mapping):
                                        is_checked = idx not in default_excluded_models
                                        checkbox = gr.Checkbox(
                                            label=f"[{idx}] {label_mapping[idx]}",
                                            value=is_checked,
                                            elem_id=f"model_checkbox_{idx}"
                                        )
                                        model_checkboxes.append((idx, checkbox))
                        
                        # Sort by index to maintain order
                        model_checkboxes.sort(key=lambda x: x[0])
                        checkbox_inputs = [cb[1] for cb in model_checkboxes]
                    
                    with gr.Column(scale=2):
                        dev_result_output = gr.HTML("", elem_id="dev_result_output_box", label="Result & Highlighted Text")
                        dev_debug_output = gr.HTML("", label="Per-Sentence Debug Information", elem_classes=["scrollable-debug"])
                
                def process_dev_analysis(text, *checkbox_values, progress=gr.Progress()):
                    """Process text with developer options and show Per-Sentence Debug Information"""
                    if not text or not text.strip():
                        return "", "", ""
                    
                    # Get excluded indices (unchecked checkboxes)
                    excluded_indices = [idx for idx, checked in enumerate(checkbox_values) if not checked]
                    
                    # Process with debug info - returns (text_out, result_msg, highlighted, debug_info)
                    text_out, result_msg, highlighted, debug_info = process_file_or_text(
                        None, text, progress, excluded_indices=excluded_indices, return_debug_info=True
                    )
                    
                    # Combine result message and highlighted text for dev_result_output
                    combined_result = f"{result_msg}<br><br><h3>Highlighted Text:</h3>{highlighted}"
                    
                    # Return combined result (for dev_result_output), highlighted text (duplicate for compatibility), and debug info (for dev_debug_output)
                    return combined_result, highlighted, debug_info if debug_info else ""
                
                # Add scan button for Developer Options
                dev_scan_btn = gr.Button("Scan", variant="primary", size="lg")
                
                # Scan button handler - shows result and highlighted text in dev_result_output, debug info in dev_debug_output
                dev_scan_btn.click(
                    process_dev_analysis,
                    inputs=[text_input] + checkbox_inputs,
                    outputs=[dev_result_output, dev_result_output, dev_debug_output]
                )
                
                # Trigger when checkboxes change - auto-scan
                for checkbox in checkbox_inputs:
                    checkbox.change(
                        process_dev_analysis,
                        inputs=[text_input] + checkbox_inputs,
                        outputs=[dev_result_output, dev_result_output, dev_debug_output]
                    )
    
    # Shared file upload handler - populates the shared text input
    def extract_text_for_tabs(file):
        """Extract text from file and populate the shared text input"""
        if file is None or file == "":
            return ""
        
        file_path = file if isinstance(file, str) else (file.name if hasattr(file, 'name') else str(file))
        if not file_path or file_path == "":
            return ""
        
        saved_path = save_uploaded_file(file_path)
        extracted_text = extract_text_from_file(saved_path)
        
        if extracted_text.startswith("Error") or extracted_text.startswith("Unsupported"):
            return ""
        
        return extracted_text

    def update_display(text_input_val, highlighted_html):
        """Update highlighted text display in AI Detection tab"""
        if highlighted_html and highlighted_html.strip():
            return gr.update(), gr.update(value=highlighted_html, visible=True)
        else:
            return gr.update(), gr.update(value="", visible=False)
    
    def disable_inputs(progress=gr.Progress()):
        """Disable text input during processing, change button to Stop"""
        # Show initial progress immediately
        if progress:
            progress(0.0, desc="Preparing...")
        # Don't disable file_upload to preserve the clear button - just return no-op update
        return (
            gr.update(interactive=False),  # Only disable text input
            gr.update(),  # File upload - no change (preserves clear button)
            gr.update(value="Stop", variant="stop", interactive=True),
            True  # Set scanning state to True
        )
    
    def enable_inputs():
        """Re-enable text input after processing, change button back to Scan"""
        return (
            gr.update(interactive=True),  # Re-enable text input
            gr.update(),  # File upload - no change (preserves clear button)
            gr.update(value="Scan", variant="primary", interactive=True),
            False  # Set scanning state to False
        )
    
    def stop_scanning():
        """Stop scanning and reset button"""
        return (
            gr.update(value="Scan", variant="primary"),
            False
        )
    
    # Process when text is changed (either manually or from file upload)
    def process_and_store(file, text_input_val, progress=gr.Progress()):
        """Process text and store results with error handling"""
        try:
            # If file is uploaded but text_input is empty, check for extraction errors
            if file and (not text_input_val or not text_input_val.strip()):
                file_path = file if isinstance(file, str) else (file.name if hasattr(file, 'name') else str(file))
                if file_path:
                    # Save uploaded file to uploads directory
                    saved_path = save_uploaded_file(file_path)
                    extracted_text = extract_text_from_file(saved_path)
                    if extracted_text.startswith("Error") or extracted_text.startswith("Unsupported"):
                        error_msg = f'<div style="color: red; font-weight: bold;">{html.escape(extracted_text)}</div>'
                        return "", error_msg, "", "", error_msg, ""
            
            text_out, result_msg, highlighted = process_file_or_text(file, text_input_val, progress)
            return text_out, result_msg, highlighted, text_out, result_msg, highlighted
        except Exception as e:
            error_msg = f'<div style="color: red; font-weight: bold;">Error: {html.escape(str(e))}</div>'
            return "", error_msg, "", "", error_msg, ""
    
    # When file is uploaded, extract text and populate text_input
    file_upload.change(
        extract_text_for_tabs, 
        inputs=[file_upload], 
        outputs=[text_input]
    )
    
    # Process all tabs function
    def process_all_tabs(file, text_input_val, grammar_lang, exclude_urls, *checkbox_values, progress=gr.Progress()):
        """Process text for all tabs"""
        # Check if text is empty
        if not text_input_val or not text_input_val.strip():
            empty_msg = '<div style="color: orange; font-weight: bold;">Please enter or upload text to analyze.</div>'
            return (
                "", empty_msg, "", "", empty_msg, "",  # AI Detection
                empty_msg,  # TOEFL
                empty_msg,  # Feedback
                empty_msg, "**Plagiarism Score:** 0.0%",  # Plagiarism
                empty_msg, "**Citations Detected:** 0",  # Citation Checker
                empty_msg, "" if not DEV_MODE else ""  # Developer Options
            )
        
        # Check language first - only English is supported
        is_valid, error_msg = check_language_and_return_error(text_input_val)
        if not is_valid:
            return (
                "", error_msg, "", "", error_msg, "",  # AI Detection
                error_msg,  # TOEFL
                error_msg,  # Feedback
                error_msg, "**Plagiarism Score:** 0.0%",  # Plagiarism
                error_msg, "**Citations Detected:** 0",  # Citation Checker
                error_msg, "" if not DEV_MODE else ""  # Developer Options
            )
        
        results = {}
        
        # AI Detection (now includes grammar checking with red underlines)
        if progress:
            progress(0.1, desc="Processing AI Detection...")
        try:
            text_out, result_msg, highlighted = process_file_or_text(file, text_input_val, progress)
            results['ai'] = (text_out, result_msg, highlighted, text_out, result_msg, highlighted)
        except Exception as e:
            error_msg = f'<div style="color: red; font-weight: bold;">Error: {html.escape(str(e))}</div>'
            results['ai'] = ("", error_msg, "", "", error_msg, "")
        
        # TOEFL Rubric
        if progress:
            progress(0.5, desc="Processing TOEFL Rubric...")
        try:
            if text_input_val and text_input_val.strip():
                scores = calculate_toefl_score(text_input_val)
                toefl_result = format_toefl_rubric(scores, text_input_val) if scores else ""
                results['toefl'] = (toefl_result,)
            else:
                results['toefl'] = ("",)
        except Exception as e:
            results['toefl'] = (f'<div style="color: red;">Error: {html.escape(str(e))}</div>',)
        
        # Writing Feedback
        if progress:
            progress(0.7, desc="Processing Writing Feedback...")
        try:
            if text_input_val and text_input_val.strip():
                feedback_result = generate_writing_feedback(text_input_val)
                results['feedback'] = (feedback_result,)
            else:
                results['feedback'] = ("",)
        except Exception as e:
            results['feedback'] = (f'<div style="color: red;">Error: {html.escape(str(e))}</div>',)
        
        # Plagiarism Check
        if progress:
            progress(0.9, desc="Processing Plagiarism Check...")
        try:
            if text_input_val and text_input_val.strip():
                exclude_list = [url.strip() for url in exclude_urls.split('\n') if url.strip()] if exclude_urls else None
                result_html, similarity = check_plagiarism(text_input_val, exclude_list)
                score_text = f"**Plagiarism Score:** {similarity:.1f}%"
                results['plagiarism'] = (result_html, score_text)
            else:
                results['plagiarism'] = ("", "**Plagiarism Score:** 0.0%")
        except Exception as e:
            results['plagiarism'] = (f'<div style="color: red;">Error: {html.escape(str(e))}</div>', "**Plagiarism Score:** 0.0%")
        
        # Citation Checker
        if progress:
            progress(0.92, desc="Detecting citations...")
        try:
            if text_input_val and text_input_val.strip():
                citation_html, citation_summary = detect_citations(text_input_val)
                results['citation'] = (citation_html, citation_summary)
            else:
                results['citation'] = ("", "**Citations Detected:** 0")
        except Exception as e:
            results['citation'] = (
                f'<div style="color: red;">Error in citation detection: {html.escape(str(e))}</div>',
                "**Citations Detected:** 0",
            )
        
        # Developer Options
        if DEV_MODE and checkbox_values and len(checkbox_values) > 0:
            if progress:
                progress(0.95, desc="Processing Developer Analysis...")
            try:
                if text_input_val and text_input_val.strip():
                    # Get excluded indices (unchecked checkboxes)
                    excluded_indices = [idx for idx, checked in enumerate(checkbox_values) if not checked]
                    
                    # Process with debug info - returns (text_out, result_msg, highlighted, debug_info)
                    text_out, result_msg, highlighted, debug_info = process_file_or_text(
                        None, text_input_val, None, excluded_indices=excluded_indices, return_debug_info=True
                    )
                    
                    # Return result message, highlighted text, and debug info separately
                    results['dev'] = (result_msg, highlighted, debug_info if debug_info else "")
                else:
                    results['dev'] = ("", "", "")
            except Exception as e:
                results['dev'] = (f'<div style="color: red;">Error: {html.escape(str(e))}</div>', "", "")
        else:
            results['dev'] = ("", "", "")
        
        if progress:
            progress(1.0, desc="All analyses complete!")
        
        # Return all results in the correct order
        return (
            results['ai'][0], results['ai'][1], results['ai'][2],  # AI Detection outputs
            results['ai'][3], results['ai'][4], results['ai'][5],  # AI Detection state
            results['toefl'][0],  # TOEFL output
            results['feedback'][0],  # Feedback output
            results['plagiarism'][0], results['plagiarism'][1],  # Plagiarism outputs
            results['citation'][0], results['citation'][1],  # Citation outputs
            results['dev'][0], results['dev'][1], results['dev'][2]  # Developer Options outputs (result, highlighted, debug_info)
        )
    
    # Process button triggers analysis for all tabs
    scan_inputs = [file_upload, text_input, grammar_language, exclude_urls_input]
    scan_outputs = [
        text_input, result_output, highlighted_text_display,  # AI Detection outputs
        last_processed_text, last_result_msg, last_highlighted_html,  # AI Detection state
        toefl_output,  # TOEFL output
        feedback_output,  # Feedback output
        plagiarism_output, plagiarism_score,  # Plagiarism outputs
        citation_output, citation_stats,  # Citation outputs
    ]
    
    # Add Developer Options inputs/outputs if DEV_MODE is enabled
    if DEV_MODE and checkbox_inputs and dev_result_output and dev_debug_output:
        scan_inputs.extend(checkbox_inputs)
        scan_outputs.extend([dev_result_output, dev_debug_output])
    
    process_btn.click(
        disable_inputs, outputs=[text_input, file_upload, process_btn, scanning_state]
    ).then(
        process_all_tabs, 
        inputs=scan_inputs, 
        outputs=scan_outputs
    ).then(
        update_display, inputs=[text_input, highlighted_text_display], outputs=[text_input, highlighted_text_display]
    ).then(
        enable_inputs, outputs=[text_input, file_upload, process_btn, scanning_state]
    ).then(
        None,
        None,
        None,
        js="""
        function() {
            // Attach grammar error click handlers after content is updated
            setTimeout(function() {
                const grammarErrors = document.querySelectorAll('.grammar-error-word, .grammar-error-underline');
                console.log('Found grammar errors:', grammarErrors.length);
                
                grammarErrors.forEach(function(errorSpan) {
                    // Remove any existing listeners
                    const newSpan = errorSpan.cloneNode(true);
                    errorSpan.parentNode.replaceChild(newSpan, errorSpan);
                    
                    // Add click handler
                    newSpan.addEventListener('click', function(e) {
                        e.stopPropagation();
                        e.preventDefault();
                        console.log('Grammar error clicked!');
                        
                        // Remove existing popups
                        const existingPopup = document.querySelector('.grammar-suggestion-popup');
                        if (existingPopup) {
                            existingPopup.remove();
                        }
                        
                        // Get error data
                        const message = newSpan.getAttribute('data-grammar-message') || 'Grammar error';
                        const suggestionsStr = newSpan.getAttribute('data-grammar-suggestions') || '';
                        const suggestions = suggestionsStr ? suggestionsStr.split('|') : [];
                        
                        console.log('Message:', message, 'Suggestions:', suggestions);
                        
                        // Create popup
                        const popup = document.createElement('div');
                        popup.className = 'grammar-suggestion-popup show';
                        
                        // Position popup
                        const rect = newSpan.getBoundingClientRect();
                        const scrollTop = window.pageYOffset || document.documentElement.scrollTop;
                        const scrollLeft = window.pageXOffset || document.documentElement.scrollLeft;
                        
                        let top = rect.top + scrollTop - 10;
                        let left = rect.left + scrollLeft;
                        
                        if (top < scrollTop) {
                            top = rect.bottom + scrollTop + 10;
                        }
                        if (left + 400 > window.innerWidth + scrollLeft) {
                            left = window.innerWidth + scrollLeft - 420;
                        }
                        if (left < scrollLeft) {
                            left = scrollLeft + 10;
                        }
                        
                        popup.style.top = top + 'px';
                        popup.style.left = left + 'px';
                        popup.style.position = 'absolute';
                        popup.style.zIndex = '10000';
                        
                        // Build popup content
                        let content = '<span class="close-btn">&times;</span>';
                        content += '<div class="message">' + message + '</div>';
                        
                        if (suggestions.length > 0) {
                            content += '<div class="suggestions"><strong>Suggestions:</strong>';
                            suggestions.forEach(function(suggestion, index) {
                                content += '<div class="suggestion-item" data-index="' + index + '">' + 
                                          suggestion + '</div>';
                            });
                            content += '</div>';
                        } else {
                            content += '<div class="suggestions">No suggestions available.</div>';
                        }
                        
                        popup.setAttribute('data-suggestions', JSON.stringify(suggestions));
                        popup.innerHTML = content;
                        document.body.appendChild(popup);
                        console.log('Popup created');
                        
                        // Handle suggestion clicks
                        popup.querySelectorAll('.suggestion-item').forEach(function(item) {
                            item.addEventListener('click', function(e) {
                                e.stopPropagation();
                                const index = parseInt(item.getAttribute('data-index') || '0');
                                const suggestions = JSON.parse(popup.getAttribute('data-suggestions') || '[]');
                                if (index >= 0 && index < suggestions.length) {
                                    const replacement = suggestions[index];
                                    if (navigator.clipboard && navigator.clipboard.writeText) {
                                        navigator.clipboard.writeText(replacement).then(function() {
                                            item.style.backgroundColor = '#4CAF50';
                                            item.style.color = 'white';
                                            setTimeout(function() {
                                                item.style.backgroundColor = '';
                                                item.style.color = '';
                                            }, 1000);
                                        }).catch(function() {
                                            alert('Suggestion: ' + replacement + '\\n\\nPlease copy this and replace the error manually.');
                                        });
                                    } else {
                                        alert('Suggestion: ' + replacement + '\\n\\nPlease copy this and replace the error manually.');
                                    }
                                }
                                setTimeout(function() {
                                    if (popup) popup.remove();
                                }, 1500);
                            });
                        });
                        
                        // Handle close button
                        popup.querySelector('.close-btn').addEventListener('click', function() {
                            popup.remove();
                        });
                        
                        // Close on outside click
                        setTimeout(function() {
                            document.addEventListener('click', function closePopup(e) {
                                if (!popup.contains(e.target) && e.target !== newSpan) {
                                    popup.remove();
                                    document.removeEventListener('click', closePopup);
                                }
                            });
            }, 100);
                    });
                });
            }, 500);
        }
        """
    )
    
    

    # with gr.Tab("AI text examples"):
    #     gr.Examples(AI_texts, inputs=text_input)
    # with gr.Tab("Human text examples"):
    #     gr.Examples(Human_texts, inputs=text_input)
    gr.Markdown(bottom_text, elem_id="bottom_text")
    
    # Inject CSS on app load using JavaScript
    def inject_css():
        return None
    
    # Escape CSS for JavaScript (pre-process outside f-string to avoid syntax error)
    escaped_css = custom_css.replace('`', '\\`').replace('$', '\\$').replace('\\', '\\\\')
    
    # Use iface.load() to inject CSS when the app loads (must be inside Blocks context)
    iface.load(
        fn=inject_css,
        js=f"""
        function() {{
            const style = document.createElement('style');
            style.textContent = `{escaped_css}`;
            if (!document.getElementById('custom-gradio-css')) {{
                style.id = 'custom-gradio-css';
                document.head.appendChild(style);
            }}
        }}
        """
    )

# Get port from environment variable or use default
port = int(os.environ.get("GRADIO_PORT", 7860))

if __name__ == "__main__":
    try:
        print(f"Starting Gradio app on port {port}...", flush=True)
        print(f"Server will be available at http://0.0.0.0:{port}", flush=True)
        iface.launch(server_name="0.0.0.0", server_port=port, share=False)
    except Exception as e:
        print(f"Error starting Gradio app: {e}", flush=True)
        import traceback
        traceback.print_exc()
        raise